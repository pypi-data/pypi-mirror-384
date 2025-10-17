# Copyright (C) 2025 dssTools Developers
# GNU General Public License v3.0+ (see LICENSE or https://www.gnu.org/licenses/gpl-3.0.txt)
"""Module for handling the batch export of ImageGenerators."""

from io import BytesIO
from pathlib import Path
from typing import Iterable, Optional, Union

import networkx as nx
from matplotlib import pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.figure import Figure
from pptx import Presentation, presentation
from pptx.util import Cm, Length

from dsstools.mapping import qualitative
from dsstools.draw import ImageGenerator, ensure_file_format
from dsstools.log.logger import get_logger

logger = get_logger("root")


class ImageCollection(list):
    """Class for exporting multiple ImageGenerators in one go."""
    def __init__(self, iterable: Optional[Iterable[ImageGenerator]]=None):
        if iterable:
            super().__init__(self._validate_image_generator(item) for item in iterable)

    def __setitem__(self, id, item):
        super().__setitem__(id, self._validate_image_generator(item))

    def append(self, item: ImageGenerator):
        """Add new ImageGenerator to ImageCollection.

        Args:
          item: Item to append to list.

        Returns:

        """
        super().append(self._validate_image_generator(item))

    def extend(self, other: Iterable[ImageGenerator]):
        """Extend existing ImageCollection with another one.

        Args:
          other: Another ImageCollection to extend with.

        Returns:

        """
        if isinstance(other, type(self)):
            super().extend(other)
        else:
            super().extend(self._validate_image_generator(item) for item in other)

    def insert(self, id, item: ImageGenerator):
        """Insert an item at a specific spot.

        Args:
          id: Spot to insert at.
          item: The item to insert.

        Returns:
          The updated ImageCollection.
        """
        super().insert(id, self._validate_image_generator(item))

    def _validate_image_generator(self, value):
        """Ensure only ImageGenerators are being added.

        Args:
          value: Argument to check

        Returns:
          value, if it was actually an ImageGenerator.

        Raises:
          TypeError, if no ImageGenerator was passed.

        """
        if isinstance(value, ImageGenerator):
            return value
        type_error = TypeError(f"ImageGenerator expected, got {type(value).__name__} instead.")
        logger.error(type_error)
        raise type_error

    def create_flipbook(self, path: Path | str, **kwargs):
        """"Creates a flipbook as PPTX or PDF depending on file ending.

        For the specific valid keyword arguments see `create_flipbook_pdf` or
        `create_flipbook_pptx` which this a wrapper for.

        Args:
            path: Path to save flipbook to. File ending decides on the internal file
                format.

        Returns:
            Either a PDF or PPTX object.
        """
        path = Path(path)
        # TODO Replace this with Henriks file ending checker
        if path.suffix == ".pdf":
            return self.create_flipbook_pdf(path, **kwargs)
        if path.suffix == ".pptx":
            return self.create_flipbook_pptx(path, **kwargs)
        value_error = ValueError(f"Unable to inpret the file ending {path.suffix}")
        logger.error(value_error)
        raise value_error


    def create_flipbook_pdf(self, path: Path | str) -> PdfPages:
        """Create PDF containing all ImageGenerators.

        Args:
          path: Path to save PDF to.

        Returns:
          Generated PDF object.
        """
        path = Path(path)
        with PdfPages(path) as pdf:
            for ig in self:
                if not ig._fig:
                    ig.draw()
                pdf.savefig(ig._fig)
            return pdf

    def create_flipbook_pptx(
        self,
        path: Path | str,
        titles: Optional[list] = None,
        left: Length = Cm(4),
        top: Length = Cm(-5.3),
        height: Length = Cm(25),
    ) -> presentation.Presentation:
        """Create PPTX containing all ImageGenerators.

        Args:
          path: Path to save file to.
          titles: Titles to give each slide. (Default value = None)
          left:  Left offset of the image on the slide, starting from upper left.
            (Default value = Cm(4))
          top: Top offset of the image on the slide, starting from upper left.
            (Default value = Cm(-5.3))
          height: Height of the image. By default uses a sensible default. If you change
            this, you might have to adapt the left and top arguments as well. (Default
            value = Cm(25))

        Returns:
          Generated PPTX object.
        """
        path = Path(path)

        if titles and len(self) != len(titles):
            value_error = ValueError(
                "Titles and image generators in the flipbook are not the same size."
            )
            logger.error(value_error)
            raise value_error


        prs = Presentation()
        # We default to 16:9 as slide layout
        prs.slide_width = 9144000
        prs.slide_height = 5144400

        for i, ig in enumerate(self):
            title_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_slide_layout)
            if titles:
                slide.shapes.title.text = titles[i]
            if not ig._fig:
                ig.draw()
            # Save the image to a byte stream to avoid disk access
            stream = BytesIO()
            ig._fig.savefig(stream, dpi=ig.dpi)
            pic = slide.shapes.add_picture(stream, left=left, top=top, height=height)

            # Move image to background
            cursor_sp = slide.shapes[0]._element
            cursor_sp.addprevious(pic._element)

        prs.save(path)
        return prs

    def create_multiple_in_one(self, fig: Figure, path: Path | str, dpi: int = 200):
        path = Path(path)

        # TODO Use Henriks file ending checker
        if len(self) != len(fig.axes):
            value_error = ValueError("List of graphs and axes in figure are not the same length.")
            logger.error(value_error)
            raise value_error

        for ig, ax in zip(self, fig.axes):
            ig.set_axis(ax).draw()
            ax.set_axis_off()

        if path:
            fpath, file_format = ensure_file_format(path,
                                                    path.suffix,
                                                    default_format=".svg")
            fig.savefig(fpath, format=file_format, dpi=dpi)

        return fig

    """
    def highlight_ego_network(self, index: int = 0, ego, path: Union[Path, str], color_ego: str = "red", color_neighbors: str = "green", color_neighbors_edges: str = "green", grey_out_rest: bool = True) -> None:
        Draws an ego-network around ego.

        Args:
            ego: a node, the ego of the network
            path: a path, the path where the image of the ego-network should be saved
            color_ego: a string, the color for the ego-node, default red
            color_neighbors: a string, the color for the neighboring nodes
            color_neighbors_edges: a string, the color for the edges between neighbors and ego
            grey_out_rest: a boolean, if True: make remaining part of graph (that is not in egonetwork) grey, if False: keep previous coloring, default true

        if index != 0:
            warnings.warn
        #create subgraph that will become the ego-network
        neighbors = list(self[index].graph.neighbors(ego))

        ego_graph = nx.ego_graph(se.graph, ego, is_directed=True)

        #create ImageGenerator for subgraph
        ig_ego_graph = ImageGenerator(ego_graph)
        ig_ego_graph.nodes.set_positions(image.nodes.positions)

        #color subgraph
        ig_ego_graph.edges.set_colors(color_neighbors_edges)
        ego_graph.nodes[ego]["ego"] = True

        for neighbor in neighbors:
            ego_graph.nodes[neighbor]["ego"] = False

        ig_ego_graph.nodes.set_colors(qualitative("ego", {True: color_ego, False: color_neighbors}))

        #color remaining nodes and edges
        if grey_out_rest:
            image.nodes.set_colors("grey")
            image.edges.set_colors("grey")

        self.append(ig_ego_graph)
        self.draw_multiple_in_one_field(path)

    def draw_multiple_in_one_field(self, path) -> None:

        Draws multiple graphs on top of each other.
        path: the path where the figure will be saved

        fig, ax = plt.subplots(1,1)

        for image in self:
            image.set_axis(ax).draw()

        fig.savefig(path)
    """
