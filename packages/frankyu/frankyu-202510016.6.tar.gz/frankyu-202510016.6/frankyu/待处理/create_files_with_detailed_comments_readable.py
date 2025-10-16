# 导入 datetime 模块，用于处理日期和时间
import datetime

# 导入 time 模块，用于记录和计算时间消耗
import time

# 导入 random 模块，用于生成随机数（此处未使用，但可扩展）
import random

# 检查并导入 PowerPoint 模块
try:
    # 首选 frankyu.ppt 模块，提供更高效的 PowerPoint 操作接口
    import frankyu.ppt as ppt_module
except ImportError:
    # 如果 frankyu.ppt 模块不可用，设置 ppt_module 为 None
    ppt_module = None
    # 使用备用的 python-pptx 模块
    from pptx import Presentation

# 检查并导入 Excel 模块
try:
    # 首选 frankyu.excel 模块，用于高效操作 Excel
    import frankyu.excel as excel_module
except ImportError:
    # 如果 frankyu.excel 模块不可用，设置 excel_module 为 None
    excel_module = None
    # 使用备用的 openpyxl 模块
    from openpyxl import Workbook

# 检查并导入 Word 模块
try:
    # 首选 pywin32 模块，用于操作 Word 应用程序
    import win32com.client as win32
except ImportError:
    # 如果 pywin32 模块不可用，设置 win32 为 None
    win32 = None
    # 使用备用的 python-docx 模块
    from docx import Document


def log_step(start_time, step_description):
    """
    记录并打印单个步骤的耗时日志。

    参数:
    - start_time (float): 步骤开始时的时间戳，用于计算耗时
    - step_description (str): 步骤描述，用于打印日志
    """
    # 计算当前耗时
    elapsed_time = time.time() - start_time

    # 打印日志信息，包含步骤描述和耗时
    print(f"[日志] {step_description} - 耗时 {elapsed_time:.2f} 秒")


def create_powerpoint(base_path="T:\\ppt\\", 
                      file_name="工作总结", 
                      title_text="这是标题", 
                      content_text="这是幻灯片内容"):
    """
    创建并保存一个 PowerPoint 演示文稿。

    参数:
    - base_path (str): 文件保存路径，默认保存到 T 盘的 ppt 文件夹
    - file_name (str): 文件名（基础部分），默认值为 "工作总结"
    - title_text (str): 幻灯片的标题内容
    - content_text (str): 幻灯片的正文内容
    """
    try:
        # 打印初始化文件路径的日志
        print("[日志] 初始化 PowerPoint 文件路径...")
        
        # 记录初始化开始时间
        start_time = time.time()

        # 获取当前时间戳，用于生成唯一文件名
        current_timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")

        # 拼接完整的 PowerPoint 文件路径
        ppt_file_path = f"{base_path}{file_name}_{current_timestamp}.pptx"

        # 打印日志，记录路径初始化完成
        log_step(start_time, "初始化文件路径完成")

        if ppt_module:
            # 如果 frankyu.ppt 模块可用
            print("[日志] 使用 frankyu.ppt 模块创建 PowerPoint 文件...")

            # 记录 PowerPoint 启动时间
            start_time = time.time()

            # 启动 PowerPoint 应用程序
            powerpoint_app = ppt_module.start_powerpoint()

            # 打印日志，记录启动完成
            log_step(start_time, "启动 PowerPoint 应用程序")

            # 创建新的 PowerPoint 演示文稿
            start_time = time.time()
            presentation = ppt_module.create_new_presentation(powerpoint_app)
            log_step(start_time, "创建 PowerPoint 演示文稿")

            # 添加幻灯片内容
            start_time = time.time()
            ppt_module.add_slide_content(
                presentation=presentation,
                title_text=title_text,
                content_text=content_text
            )
            log_step(start_time, "添加幻灯片内容")

            # 保存 PowerPoint 文件
            start_time = time.time()
            presentation.SaveAs(ppt_file_path)
            log_step(start_time, "保存 PowerPoint 文件")
        else:
            # 如果 frankyu.ppt 模块不可用，改用 python-pptx
            print("[日志] 使用 python-pptx 创建 PowerPoint 文件...")

            # 开始创建演示文稿
            start_time = time.time()

            # 使用 python-pptx 创建演示文稿对象
            presentation = Presentation()

            # 获取幻灯片的默认布局
            slide_layout = presentation.slide_layouts[1]

            # 添加一张幻灯片
            slide = presentation.slides.add_slide(slide_layout)

            # 设置幻灯片标题
            slide.shapes.title.text = title_text

            # 设置幻灯片正文内容
            slide.placeholders[1].text = content_text

            # 保存演示文稿
            presentation.save(ppt_file_path)

            # 打印保存完成日志
            log_step(start_time, "使用 python-pptx 保存 PowerPoint 文件")

        # 打印 PowerPoint 文件保存成功的日志
        print(f"[日志] PowerPoint 文件保存成功: {ppt_file_path}")

        # 返回保存的文件路径
        return ppt_file_path
    except Exception as e:
        # 如果过程中发生异常，打印错误信息
        print(f"[错误] 创建 PowerPoint 文件失败: {e}")

        # 返回 None 表示失败
        return None



def create_powerpoint(base_path="T:\\ppt\\", 
                      file_name="工作总结", 
                      title_text="这是标题", 
                      content_text="这是幻灯片内容"):
    """
    创建并保存一个 PowerPoint 演示文稿。

    参数:
    - base_path (str): 文件保存路径，默认保存到 T 盘的 ppt 文件夹
    - file_name (str): 文件名（基础部分），默认值为 "工作总结"
    - title_text (str): 幻灯片的标题内容
    - content_text (str): 幻灯片的正文内容

    返回:
    - dict: 包含文件路径、PPT对象和PPT程序对象的词典
    import os
								
    """
    import os
    try:
        # 打印初始化文件路径的日志
        print("[日志] 初始化 PowerPoint 文件路径...")

        # 记录初始化开始时间
        start_time = time.time()

        # 获取当前时间戳，用于生成唯一文件名
        current_timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")

        # 拼接完整的 PowerPoint 文件路径
        ppt_file_path = os.path.join(base_path, f"{file_name}_{current_timestamp}.pptx")

        # 打印日志，记录路径初始化完成
        log_step(start_time, "初始化文件路径完成")

        powerpoint_app = None
        presentation = None

        if ppt_module:
            # 如果 frankyu.ppt 模块可用
            print("[日志] 使用 frankyu.ppt 模块创建 PowerPoint 文件...")

            # 记录 PowerPoint 启动时间
            start_time = time.time()

            # 启动 PowerPoint 应用程序
            powerpoint_app = ppt_module.start_powerpoint()

            # 打印日志，记录启动完成
            log_step(start_time, "启动 PowerPoint 应用程序")

            # 创建新的 PowerPoint 演示文稿
            start_time = time.time()
            presentation = ppt_module.create_new_presentation(powerpoint_app)
            log_step(start_time, "创建 PowerPoint 演示文稿")

            # 添加幻灯片内容
            start_time = time.time()
            ppt_module.add_slide_content(
                presentation=presentation,
                title_text=title_text,
                content_text=content_text
            )
            log_step(start_time, "添加幻灯片内容")

            # 保存 PowerPoint 文件
            start_time = time.time()
            presentation.SaveAs(ppt_file_path)
            log_step(start_time, "保存 PowerPoint 文件")
        else:
            # 如果 frankyu.ppt 模块不可用，改用 python-pptx
            print("[日志] 使用 python-pptx 创建 PowerPoint 文件...")

            # 开始创建演示文稿
            start_time = time.time()

            # 使用 python-pptx 创建演示文稿对象
            presentation = Presentation()

            # 获取幻灯片的默认布局
            slide_layout = presentation.slide_layouts[1]

            # 添加一张幻灯片
            slide = presentation.slides.add_slide(slide_layout)

            # 设置幻灯片标题
            slide.shapes.title.text = title_text

            # 设置幻灯片正文内容
            slide.placeholders[1].text = content_text

            # 保存演示文稿
            presentation.save(ppt_file_path)

            # 打印保存完成日志
            log_step(start_time, "使用 python-pptx 保存 PowerPoint 文件")

        # 打印 PowerPoint 文件保存成功的日志
        print(f"[日志] PowerPoint 文件保存成功: {ppt_file_path}")

        # 返回词典，包括文件路径、PPT对象和PPT程序对象
        return {
            "file_path": ppt_file_path,
            "presentation": presentation,
            "powerpoint_app": powerpoint_app
        }
    except Exception as e:
        # 如果过程中发生异常，打印错误信息
        print(f"[错误] 创建 PowerPoint 文件失败: {e}")

        # 返回 None 表示失败
        return None





# 示例调用
if __name__ == "__main__":
    try:
        # 调用 create_powerpoint 函数，创建 PowerPoint 文件
        ppt_path = create_powerpoint()

        # 如果文件创建成功，打印文件路径
        if ppt_path:
            print(f"[日志] 生成的 PowerPoint 文件路径: {ppt_path}")
    except Exception as e:
        # 捕获可能的异常并打印错误日志
        print(f"[未处理的错误] 创建 PowerPoint 文件失败: {e}")

    print("over")  # 程序结束标志