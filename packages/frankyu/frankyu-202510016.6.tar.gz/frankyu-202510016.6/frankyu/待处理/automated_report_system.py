# -*- coding: utf-8 -*-
"""
自动化办公报告生成系统（优先使用 frankyu 和 pywin32/xlwings，提供模块回退功能）
Automated Office Report Generation System 
(Prioritize frankyu and pywin32/xlwings with fallback)
"""

import datetime
import os
import traceback
from typing import Tuple, Optional

# 动态导入模块，用于优先加载 frankyu 或其他模块
import importlib

# ==================== 模块加载器 ====================
class ModuleLoader:
    @staticmethod
    def load(
        primary: str,
        fallback: str,
        required_functions: list
    ) -> Optional[object]:
        """
        优先加载主模块，失败时加载备用模块
        Priority loading with fallback
        
        Args:
            primary: 首选模块路径 (e.g. "frankyu.ppt")
            fallback: 备用模块路径 (e.g. "backup.ppt")
            required_functions: 模块必须包含的函数列表
            
        Returns:
            加载的模块对象 (或 None 表示加载失败)
        """
        for module_name in [primary, fallback]:
            try:
                print(f"[INFO] 尝试加载模块: {module_name}")
                module = importlib.import_module(module_name)

                # 验证模块是否包含所需函数
                missing_funcs = [
                    func for func in required_functions
                    if not hasattr(module, func)
                ]
                if missing_funcs:
                    print(
                        f"[WARNING] 模块 {module_name} 缺少必要函数: {missing_funcs}"
                    )
                    continue

                print(f"[SUCCESS] 模块加载成功: {module_name}")
                return module

            except ImportError as e:
                print(
                    f"[WARNING] 模块加载失败 {module_name}: {str(e)}"
                )
                continue
            except Exception as e:
                print(
                    f"[ERROR] 模块验证失败 {module_name}: {str(e)}"
                )
                continue

        print(
            f"[INFO] 无法加载模块 {primary} 或 {fallback}，将尝试其他方案"
        )
        return None

# ==================== 核心功能 ====================
def generate_timestamp(
    time_format: str = "%Y%m%d%H%M%S"
) -> str:
    """
    生成标准化时间戳
    Generate standardized timestamp
    
    Args:
        time_format: 时间格式字符串 (默认"%Y%m%d%H%M%S")
        
    Returns:
        格式化后的时间字符串
    """
    try:
        return datetime.datetime.now().strftime(time_format)
    except Exception as e:
        print(f"[ERROR] 时间戳生成失败: {str(e)}")
        print(traceback.format_exc())
        return f"backup_{int(datetime.datetime.now().timestamp())}"

def ensure_directory(
    path: str = os.getcwd(),
    is_file_path: bool = False
) -> str:
    """
    确保目录存在，不存在则自动创建
    Ensure directory exists, create if not
    
    Args:
        path: 目标路径
        is_file_path: 是否为文件路径
        
    Returns:
        标准化后的目录路径（以分隔符结尾）
    """
    try:
        dir_path = (
            os.path.dirname(path) if is_file_path else path
        )
        dir_path = os.path.abspath(dir_path)

        if not os.path.exists(dir_path):
            os.makedirs(dir_path, exist_ok=True)
            print(f"[INFO] 目录创建成功: {dir_path}")

        return os.path.join(dir_path, "")
    except Exception as e:
        print(f"[ERROR] 目录处理失败: {str(e)}")
        print(f"[INFO] 使用当前目录: {os.getcwd()}")
        return os.path.join(os.getcwd(), "")

# ==================== 文件操作 ====================
def create_ppt_with_frankyu_or_others(
    report_name: str,
    timestamp: str,
    title: str,
    content: str,
    save_dir: str
) -> Optional[object]:
    """
    优先使用 frankyu 和 pywin32 创建 PPT，
    如果不可用则回退到 python-pptx
    """
    try:
        # 尝试加载 frankyu 模块
        ppt_module = ModuleLoader.load(
            primary="frankyu.ppt",
            fallback="backup.ppt",
            required_functions=[
                "start_powerpoint",
                "create_new_presentation",
                "add_slide_content"
            ]
        )
        if ppt_module:
            app = ppt_module.start_powerpoint()
            ppt_file = ppt_module.create_new_presentation(app)
            ppt_module.add_slide_content(
                presentation=ppt_file,
                title_text=title,
                content_text=content
            )
            save_path = os.path.join(
                ensure_directory(save_dir),
                f"{report_name}_{timestamp}.pptx"
            )
            ppt_file.SaveAs(save_path)
            print(f"[SUCCESS] 使用 frankyu 创建 PPT 成功: {save_path}")
            return ppt_file

    except Exception as e:
        print(f"[ERROR] 使用 frankyu 创建 PPT 失败: {str(e)}")
        print(traceback.format_exc())

    # 尝试使用 pywin32 创建 PPT
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        presentation = powerpoint.Presentations.Add()
        slide = presentation.Slides.Add(1, 1)
        slide.Shapes[0].TextFrame.TextRange.Text = title
        slide.Shapes[1].TextFrame.TextRange.Text = content
        save_path = os.path.join(
            ensure_directory(save_dir),
            f"{report_name}_{timestamp}.pptx"
        )
        presentation.SaveAs(save_path)
        print(f"[SUCCESS] 使用 pywin32 创建 PPT 成功: {save_path}")
        return presentation
    except ImportError:
        print("[WARNING] pywin32 不可用，尝试回退到 python-pptx")
    except Exception as e:
        print(f"[ERROR] 使用 pywin32 创建 PPT 失败: {str(e)}")
        print(traceback.format_exc())

    # 回退到 python-pptx
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        save_path = os.path.join(
            ensure_directory(save_dir),
            f"{report_name}_{timestamp}.pptx"
        )
        prs.save(save_path)
        print(f"[SUCCESS] 使用 python-pptx 创建 PPT 成功: {save_path}")
        return prs
    except Exception as e:
        print(f"[ERROR] 使用 python-pptx 创建 PPT 失败: {str(e)}")
        print(traceback.format_exc())
        return None

def create_excel_with_xlwings_or_others(
    report_name: str,
    timestamp: str,
    save_dir: str
) -> Optional[object]:
    """
    优先使用 frankyu 和 xlwings 创建 Excel，
    如果不可用则回退到 openpyxl
    """
    try:
        excel_module = ModuleLoader.load(
            primary="frankyu.excel",
            fallback="backup.excel",
            required_functions=[
                "initialize_excel",
                "create_workbook",
                "save_workbook"
            ]
        )
        if excel_module:
            app = excel_module.initialize_excel()
            workbook = excel_module.create_workbook(app)
            save_path = os.path.join(
                ensure_directory(save_dir),
                f"{report_name}_{timestamp}.xlsx"
            )
            excel_module.save_workbook(
                workbook=workbook,
                file_path=save_path
            )
            print(f"[SUCCESS] 使用 frankyu 创建 Excel 成功: {save_path}")
            return workbook

    except Exception as e:
        print(f"[ERROR] 使用 frankyu 创建 Excel 失败: {str(e)}")
        print(traceback.format_exc())

    # 尝试使用 xlwings 创建 Excel
    try:
        import xlwings as xw
        app = xw.App(visible=True)
        workbook = app.books.add()
        sheet = workbook.sheets[0]
        sheet.name = "工作报告"
        sheet.range("A1").value = ["项目", "进度", "备注"]
        sheet.range("A2").value = ["示例项目", "完成80%", "无"]
        save_path = os.path.join(
            ensure_directory(save_dir),
            f"{report_name}_{timestamp}.xlsx"
        )
        workbook.save(save_path)
        print(f"[SUCCESS] 使用 xlwings 创建 Excel 成功: {save_path}")
        return workbook
    except ImportError:
        print("[WARNING] xlwings 不可用，尝试回退到 openpyxl")
    except Exception as e:
        print(f"[ERROR] 使用 xlwings 创建 Excel 失败: {str(e)}")
        print(traceback.format_exc())

    # 回退到 openpyxl
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "工作报告"
        ws.append(["项目", "进度", "备注"])
        ws.append(["示例项目", "完成80%", "无"])
        save_path = os.path.join(
            ensure_directory(save_dir),
            f"{report_name}_{timestamp}.xlsx"
        )
        wb.save(save_path)
        print(f"[SUCCESS] 使用 openpyxl 创建 Excel 成功: {save_path}")
        return wb
    except Exception as e:
        print(f"[ERROR] 使用 openpyxl 创建 Excel 失败: {str(e)}")
        print(traceback.format_exc())
        return None