import streamlit as st
import tempfile
import os
from PIL import Image # Pillow 函式庫，用於圖片處理
import shutil # 用於刪除非空目錄
import time # 用於延遲，確保文件操作完成

# --- 嘗試導入各函式庫，並追蹤其可用性 ---
# 建立一個字典來追蹤每個函式庫的載入狀態
libraries_available = {
    "pandas": False,
    "PyPDF2": False,
    "docx": False, # 對應 python-docx
    "pptx": False, # 對應 python-pptx
    "openpyxl": False,
    "pywin32": False # 用於舊版 .doc 和 .ppt 檔案處理
}

# 嘗試導入 pandas
try:
    import pandas as pd
    libraries_available["pandas"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `pandas` 函式庫。\n"
               "Excel 檔案 (`.xls`, `.xlsx`) 處理功能將無法使用。\n"
               "請運行 `pip install pandas`。")

# 嘗試導入 PyPDF2
try:
    import PyPDF2
    libraries_available["PyPDF2"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `PyPDF2` 函式庫。\n"
               "PDF 檔案處理功能將無法使用。\n"
               "請運行 `pip install PyPDF2`。")

# 嘗試導入 python-docx
try:
    from docx import Document
    libraries_available["docx"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `python-docx` 函式庫。\n"
               "Word (.docx) 檔案處理功能將無法使用。\n"
               "請運行 `pip install python-docx`。")

# 嘗試導入 python-pptx
try:
    from pptx import Presentation
    libraries_available["pptx"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `python-pptx` 函式庫。\n"
               "PowerPoint (.pptx) 檔案處理功能將無法使用。\n"
               "請運行 `pip install python-pptx`。")

# 嘗試導入 openpyxl
# openpyxl 通常與 pandas 一起安裝，但為確保健壯性，仍單獨檢查
if libraries_available["pandas"]:
    try:
        import openpyxl
        libraries_available["openpyxl"] = True
    except ImportError:
        st.warning("⚠️ **警告：** 缺少 `openpyxl` 函式庫。\n"
                   "Excel 的高級功能（如圖片提取）可能受限。\n"
                   "請運行 `pip install openpyxl`。")
else:
    libraries_available["openpyxl"] = False

# 嘗試導入 pywin32 (用於 .doc 和 .ppt 檔案轉換)
try:
    import win32com.client
    import pythoncom # 導入 pythoncom 用於 CoInitialize/CoUninitialize
    libraries_available["pywin32"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `pywin32` 函式庫。\n"
               "舊版 Word (.doc) 及 PowerPoint (.ppt) 檔案處理功能將無法使用。\n"
               "此功能**僅限 Windows 系統並需安裝 Microsoft Office**。\n"
               "請運行 `pip install pywin32`。")
except Exception as e:
    st.warning(f"⚠️ **警告：** 導入 `pywin32` 時發生錯誤: {e}\n"
               "舊版 Word (.doc) 及 PowerPoint (.ppt) 檔案處理功能將無法使用。\n"
               "此功能**僅限 Windows 系統並需安裝 Microsoft Office**。")


# Pillow 函式庫 (PIL) 是基礎圖片處理功能，如果它都無法載入，則直接停止應用程式
try:
    # Image 已經在頂部導入，如果上面沒有報錯，說明 PIL 已經可用
    pass
except ImportError:
    st.error("❌ **嚴重錯誤：** 缺少 `Pillow` 函式庫。\n"
             "圖片顯示與處理功能將無法使用。\n"
             "請運行 `pip install Pillow`。")
    st.stop() # 如果最核心的圖片庫都缺失，則無法繼續運行

# --- Streamlit 應用程式介面 ---
st.set_page_config(layout="wide") # 設置頁面佈局為寬模式，更好顯示表格
st.title("通用文件上傳與處理工具")
st.write("請上傳圖片、音訊、視訊、Excel、PDF、Word、PowerPoint 或文字檔案，\n"
         "我會嘗試處理並顯示其內容。")

uploaded_file = st.file_uploader("選擇一個檔案上傳", type=[
    "jpg", "jpeg", "png", "gif",            # 圖片
    "mp3", "wav", "ogg",                    # 音訊
    "mp4", "mov", "avi",                    # 視訊
    "xls", "xlsx",                          # Excel
    "pdf",                                  # PDF
    "doc", "docx",                          # Word (注意：.doc 僅限 Windows + Office)
    "ppt", "pptx",                          # PowerPoint (注意：.ppt 僅限 Windows + Office)
    "txt"                                   # 文字檔案
])

# 初始化 temp_file_path 和轉換後的臨時路徑，確保它們總是有定義，以便最後進行清理
temp_file_path = None
doc_temp_docx_path = None # 用於 .doc 轉換成 .docx 的臨時路徑
ppt_temp_pptx_path = None # 用於 .ppt 轉換成 .pptx 的臨時路徑

if uploaded_file:
    # 取得檔案的 MIME 類型
    file_type = uploaded_file.type
    st.info(f"檢測到的檔案類型: **{file_type}**")

    # --- 將上傳檔案儲存到暫存位置 ---
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name # 在這裡指定路徑

    st.success(f"檔案已暫存至: `{temp_file_path}`")

    # --- 根據檔案類型處理檔案 ---

    # 圖片檔案
    if file_type.startswith("image/"):
        st.subheader("🖼️ 圖片檔案處理")
        try:
            image = Image.open(temp_file_path)
            st.image(image,
                     caption=f"上傳圖片: {uploaded_file.name}",
                     use_container_width=True)
            st.write(f"圖片尺寸: **{image.size[0]} x {image.size[1]} 像素**")
            st.write(f"圖片格式: **{image.format}**")
            image.close() # 明確關閉圖片檔案句柄
        except Exception as e:
            st.error(f"處理圖片檔案時發生錯誤: {e}")

    # 音訊檔案
    elif file_type.startswith("audio/"):
        st.subheader("🎵 音訊檔案處理")
        st.audio(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放音訊。\n"
                "若需更複雜的音訊處理，需使用專用 Python 函式庫。")

    # 視訊檔案
    elif file_type.startswith("video/"):
        st.subheader("🎬 視訊檔案處理")
        st.video(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放視訊。\n"
                "若需視訊分析或編輯，需使用如 OpenCV 等函式庫。")

    # Excel 檔案 (xls, xlsx)
    elif file_type == "application/vnd.ms-excel" or \
         file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        st.subheader("📊 Excel 檔案處理")
        if libraries_available["pandas"]: # 檢查 pandas 是否可用
            excel_file_handle = None
            temp_image_dir = None # 初始化臨時圖片目錄變數
            try:
                excel_file_handle = pd.ExcelFile(temp_file_path)
                sheet_names = excel_file_handle.sheet_names

                st.success("Excel 檔案讀取成功！")
                st.write(f"此 Excel 檔案包含 **{len(sheet_names)} 個工作表**。")

                selected_sheet = st.selectbox("選擇要查看的工作表:",
                                              sheet_names,
                                              index=0,
                                              key="excel_sheet_select")

                df_selected = pd.read_excel(excel_file_handle,
                                            sheet_name=selected_sheet)

                st.write(f"**顯示工作表: `{selected_sheet}` 的內容 (預覽前 100 行)**")
                st.dataframe(df_selected.head(100), use_container_width=True)

                if len(df_selected) > 100:
                    if st.checkbox(f"點擊查看 `{selected_sheet}` 的完整表格\n"
                                   "(可能需要更多記憶體和時間)",
                                   key=f"show_full_df_{selected_sheet}"):
                        st.dataframe(df_selected, use_container_width=True)

                st.subheader(f"📝 工作表: `{selected_sheet}` 中的文字內容")
                text_content = []
                for col in df_selected.columns:
                    for cell_value in df_selected[col].astype(str):
                        if pd.notna(cell_value) and len(cell_value.strip()) > 0:
                            text_content.append(cell_value.strip())

                if text_content:
                    with st.expander("點擊查看所有提取的儲存格文字"):
                        st.text_area("提取的文字:",
                                     "\n---\n".join(text_content),
                                     height=300)
                else:
                    st.info("此工作表的儲存格中未找到任何顯著的文字內容。")

                if file_extension == ".xlsx":
                    st.subheader(f"🖼️ 工作表: `{selected_sheet}` 中的圖片")
                    if libraries_available["openpyxl"]: # 檢查 openpyxl 是否可用
                        image_found = False
                        try:
                            temp_image_dir = tempfile.mkdtemp()
                            workbook = openpyxl.load_workbook(temp_file_path)
                            sheet = workbook[selected_sheet]

                            if hasattr(sheet, '_images') and sheet._images:
                                total_images = len(sheet._images)
                                max_images_to_preview = 5

                                st.write(f"在 `{selected_sheet}` 中找到 **{total_images}** 張圖片。")

                                if total_images > max_images_to_preview:
                                    st.info(f"為了效能，**僅預覽前 {max_images_to_preview} 張圖片**。")
                                    if st.checkbox("顯示所有圖片 (可能導致卡頓，請謹慎)",
                                                   key=f"show_all_images_{selected_sheet}"):
                                        max_images_to_preview = total_images

                                for i, img in enumerate(sheet._images):
                                    if i >= max_images_to_preview:
                                        break

                                    image_ext = ".png"
                                    if hasattr(img, 'ext') and img.ext:
                                        image_ext = f".{img.ext}"
                                    elif hasattr(img.ref, 'embed'):
                                        img_data = img.ref.embed
                                        if img_data and hasattr(img_data, 'mime_type') and img_data.mime_type:
                                            mime_type = img_data.mime_type
                                            if 'image/png' in mime_type:
                                                image_ext = '.png'
                                            elif 'image/jpeg' in mime_type or 'image/jpg' in mime_type:
                                                image_ext = '.jpg'
                                            # 您可以根據需要添加更多 mime 類型

                                    image_filename = f"image_{i}{image_ext}"
                                    image_path = os.path.join(temp_image_dir, image_filename)

                                    try:
                                        pil_image = Image.open(img.ref)
                                        if pil_image.mode in ('RGBA', 'P'):
                                            pil_image = pil_image.convert('RGB')
                                        pil_image.save(image_path, format='JPEG', quality=85)

                                        st.image(image_path,
                                                 caption=f"圖片 {i+1} 來自 {selected_sheet}",
                                                 use_container_width=True)
                                        image_found = True
                                    except Exception as img_err:
                                        st.warning(f"無法儲存圖片 {i+1} 到 {image_path}。\n"
                                                   f"錯誤: {img_err}\n"
                                                   "這可能是因為圖片格式不被 PIL 支援，或者圖片數據有問題。")

                            if not image_found:
                                st.info("此工作表中未找到任何嵌入圖片。")

                        except Exception as e:
                            st.error(f"從 Excel 檔案中提取圖片時發生錯誤\n"
                                     f"(圖片提取僅支援 .xlsx): {e}")
                        finally:
                            if temp_image_dir and os.path.exists(temp_image_dir):
                                try:
                                    shutil.rmtree(temp_image_dir)
                                except Exception as e:
                                    st.warning(f"無法刪除臨時圖片目錄 {temp_image_dir}。\n"
                                               f"原因: {e}")
                    else:
                        st.warning("`openpyxl` 函式庫未載入，無法進行 Excel 圖片提取。\n"
                                   "請檢查之前的警告訊息以了解如何安裝。")
                else:
                    st.info("圖片提取僅支援 .xlsx 檔案。")

            except Exception as e:
                st.error(f"處理 Excel 檔案時發生錯誤: {e}")
            finally:
                if excel_file_handle:
                    excel_file_handle.close()
        else:
            st.warning("`pandas` 函式庫未載入，無法處理 Excel 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # PDF 檔案
    elif file_type == "application/pdf":
        st.subheader("📄 PDF 檔案處理")
        if libraries_available["PyPDF2"]: # 檢查 PyPDF2 是否可用
            pdf_file = None
            try:
                pdf_file = open(temp_file_path, "rb")
                reader = PyPDF2.PdfReader(pdf_file)
                num_pages = len(reader.pages)
                st.write(f"PDF 檔案總頁數: **{num_pages}**")

                extracted_text = []
                for i in range(min(num_pages, 3)): # 最多提取前 3 頁的文字
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        extracted_text.append(f"--- 第 {i+1} 頁 ---\n"
                                              f"{text[:500]}..." if len(text) > 500 else
                                              f"--- 第 {i+1} 頁 ---\n{text}")
                    else:
                        extracted_text.append(f"--- 第 {i+1} 頁 ---\n"
                                              "(無法提取文字，可能為掃描件或圖像式 PDF)")

                if extracted_text:
                    st.text_area("部分文字內容預覽:",
                                 "\n\n".join(extracted_text),
                                 height=300)
                else:
                    st.info("無法從 PDF 中提取任何文字內容。")

            except Exception as e:
                st.error(f"處理 PDF 檔案時發生錯誤: {e}")
                st.info("注意: `PyPDF2` 主要用於文字型 PDF，\n"
                        "可能無法從掃描或圖像型 PDF 中提取文字。")
            finally:
                if pdf_file:
                    pdf_file.close()
        else:
            st.warning("`PyPDF2` 函式庫未載入，無法處理 PDF 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # Word 檔案 (.doc 和 .docx)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or \
         file_type == "application/msword":
        st.subheader("📝 Word 檔案處理")
        actual_file_path = temp_file_path # 預設使用原始暫存路徑

        # 如果是舊版 .doc 檔案，嘗試轉換為 .docx
        if file_type == "application/msword":
            if libraries_available["pywin32"] and libraries_available["docx"]:
                st.info("檢測到舊版 Word (.doc) 檔案，嘗試在後台轉換為 .docx 格式...")
                word = None # 初始化 Word 應用程式對象
                try:
                    # 嘗試初始化 COM
                    pythoncom.CoInitialize()
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False # 不顯示 Word 應用程式介面
                    doc = word.Documents.Open(temp_file_path)

                    doc_temp_docx_path = tempfile.mktemp(suffix=".docx")
                    doc.SaveAs(doc_temp_docx_path, FileFormat=16) # FileFormat=16 表示 .docx
                    doc.Close()
                    # word.Quit() 放在 finally 確保執行
                    actual_file_path = doc_temp_docx_path
                    st.success("檔案已成功轉換為 .docx 格式。")
                except Exception as e:
                    st.error(f"轉換舊版 Word (.doc) 檔案時發生錯誤: {e}\n"
                             "請確保您的 Windows 系統已安裝 Microsoft Office。")
                    st.warning("無法處理此 .doc 檔案。")
                    actual_file_path = None # 設置為 None，表示無法處理
                finally:
                    # 無論成功失敗，都嘗試退出 Word 應用程式並清理 COM
                    if word:
                        word.Quit()
                    pythoncom.CoUninitialize() # 清理 COM
            else:
                st.warning("處理舊版 Word (.doc) 檔案需要 `pywin32` 和 `python-docx` 函式庫，\n"
                           "且僅限 Windows 系統並安裝 Microsoft Office。\n"
                           "請檢查之前的警告訊息以了解如何安裝。")
                actual_file_path = None # 無法處理

        if actual_file_path and libraries_available["docx"]:
            try:
                document = Document(actual_file_path)
                full_text = []
                for para in document.paragraphs:
                    full_text.append(para.text)

                doc_text = "\n".join(full_text)
                if doc_text:
                    st.write(f"文件總段落數: **{len(document.paragraphs)}**")
                    st.text_area("文件內容預覽 (部分):",
                                 doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text,
                                 height=300)
                else:
                    st.info("無法從 Word 檔案中提取任何文字內容。")
            except Exception as e:
                st.error(f"處理 Word 檔案時發生錯誤: {e}")
                st.warning("這可能是檔案損壞或格式問題。")
        elif actual_file_path is None:
             # 已經在 .doc 轉換失敗時給出了警告
             pass
        else: # 如果 docx 函式庫本身就沒載入
            st.warning("`python-docx` 函式庫未載入，無法處理 Word 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # PowerPoint 檔案 (.ppt 和 .pptx)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or \
         file_type == "application/vnd.ms-powerpoint": # MIME for .ppt
        st.subheader("💡 PowerPoint 檔案處理")
        actual_file_path = temp_file_path # 預設使用原始暫存路徑

        # 如果是舊版 .ppt 檔案，嘗試轉換為 .pptx
        if file_type == "application/vnd.ms-powerpoint":
            if libraries_available["pywin32"] and libraries_available["pptx"]:
                st.info("檢測到舊版 PowerPoint (.ppt) 檔案，嘗試在後台轉換為 .pptx 格式...")
                ppt_app = None # 初始化 PowerPoint 應用程式對象
                try:
                    # 嘗試初始化 COM
                    pythoncom.CoInitialize()
                    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                    ppt_app.Visible = False # 不顯示 PowerPoint 應用程式介面
                    presentation = ppt_app.Presentations.Open(temp_file_path, WithWindow=False)

                    ppt_temp_pptx_path = tempfile.mktemp(suffix=".pptx")
                    presentation.SaveAs(ppt_temp_pptx_path, FileFormat=24) # 24 for ppOpenXMLPresentation
                    presentation.Close()
                    # ppt_app.Quit() 放在 finally 確保執行
                    actual_file_path = ppt_temp_pptx_path
                    st.success("檔案已成功轉換為 .pptx 格式。")
                except Exception as e:
                    st.error(f"轉換舊版 PowerPoint (.ppt) 檔案時發生錯誤: {e}\n"
                             "請確保您的 Windows 系統已安裝 Microsoft Office。")
                    st.warning("無法處理此 .ppt 檔案。")
                    actual_file_path = None # 設置為 None，表示無法處理
                finally:
                    # 無論成功失敗，都嘗試退出 PowerPoint 應用程式並清理 COM
                    if ppt_app:
                        ppt_app.Quit()
                    pythoncom.CoUninitialize() # 清理 COM
            else:
                st.warning("處理舊版 PowerPoint (.ppt) 檔案需要 `pywin32` 和 `python-pptx` 函式庫，\n"
                           "且僅限 Windows 系統並安裝 Microsoft Office。\n"
                           "請檢查之前的警告訊息以了解如何安裝。")
                actual_file_path = None # 無法處理

        if actual_file_path and libraries_available["pptx"]:
            try:
                prs = Presentation(actual_file_path)
                total_slides = len(prs.slides)
                st.write(f"簡報總頁數: **{total_slides}**")

                presentation_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        presentation_text.append(f"--- 投影片 {i+1} ---\n" + "\n".join(slide_text))

                if presentation_text:
                    st.text_area("部分簡報內容預覽:",
                                 "\n\n".join(presentation_text[:min(total_slides, 3)]),
                                 height=300)
                else:
                    st.info("無法從 PowerPoint 檔案中提取任何文字內容。")
            except Exception as e:
                st.error(f"處理 PowerPoint 檔案時發生錯誤: {e}")
                st.warning("這可能是檔案損壞或格式問題。")
        elif actual_file_path is None:
             # 已經在 .ppt 轉換失敗時給出了警告
             pass
        else: # 如果 pptx 函式庫本身就沒載入
            st.warning("`python-pptx` 函式庫未載入，無法處理 PowerPoint 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # 文字檔案
    elif file_type.startswith("text/"):
        st.subheader("📜 文字檔案處理")
        text_file = None
        try:
            text_file = open(temp_file_path, "r", encoding="utf-8")
            content = text_file.read()
            st.write(f"檔案大小: **{len(content.encode('utf-8')) / 1024:.2f} KB**")
            st.text_area("檔案內容預覽:",
                         content[:2000] + "..." if len(content) > 2000 else content,
                         height=300)
        except UnicodeDecodeError:
            st.warning("嘗試使用 `latin-1` 編碼讀取文字檔案...")
            try:
                text_file = open(temp_file_path, "r", encoding="latin-1")
                content = text_file.read()
                st.text_area("檔案內容預覽:",
                             content[:2000] + "..." if len(content) > 2000 else content,
                             height=300)
            except Exception as e:
                st.error(f"處理文字檔案時發生編碼錯誤: {e}")
        except Exception as e:
            st.error(f"處理文字檔案時發生錯誤: {e}")
        finally:
            if text_file:
                text_file.close()

    # 其他不支援的檔案類型
    else:
        st.warning(f"🤔 抱歉，目前不支援處理類型為 `{file_type}` 的檔案。")

    # --- 清理暫存檔案 ---
    # 原始上傳檔案的暫存檔
    if temp_file_path and os.path.exists(temp_file_path):
        try:
            os.unlink(temp_file_path)
            st.info("原始上傳的暫存檔案已成功刪除。")
        except PermissionError:
            st.warning("⚠️ **權限錯誤：** 無法刪除原始暫存檔案。\n"
                       "檔案可能仍被系統或其他程式使用中。")
        except Exception as e:
            st.error(f"刪除原始暫存檔案時發生未知錯誤: {e}")

    # 如果有 .doc 轉換的臨時 .docx 檔案，也清理掉
    if doc_temp_docx_path and os.path.exists(doc_temp_docx_path):
        # 增加短暫延遲，確保Word應用程式完全關閉釋放文件句柄
        time.sleep(0.5)
        try:
            os.unlink(doc_temp_docx_path)
            st.info("舊版 Word 轉換後的臨時 .docx 檔案已成功刪除。")
        except PermissionError:
            st.warning("⚠️ **權限錯誤：** 無法刪除舊版 Word 轉換後的臨時 .docx 檔案。\n"
                       "檔案可能仍被系統或其他程式使用中。")
        except Exception as e:
            st.error(f"刪除舊版 Word 轉換後的臨時 .docx 檔案時發生未知錯誤: {e}")

    # 如果有 .ppt 轉換的臨時 .pptx 檔案，也清理掉
    if ppt_temp_pptx_path and os.path.exists(ppt_temp_pptx_path):
        # 增加短暫延遲，確保PowerPoint應用程式完全關閉釋放文件句柄
        time.sleep(0.5)
        try:
            os.unlink(ppt_temp_pptx_path)
            st.info("舊版 PowerPoint 轉換後的臨時 .pptx 檔案已成功刪除。")
        except PermissionError:
            st.warning("⚠️ **權限錯誤：** 無法刪除舊版 PowerPoint 轉換後的臨時 .pptx 檔案。\n"
                       "檔案可能仍被系統或其他程式使用中。")
        except Exception as e:
            st.error(f"刪除舊版 PowerPoint 轉換後的臨時 .pptx 檔案時發生未知錯誤: {e}")

else:
    pass # 沒有上傳檔案時不執行任何操作