import os
import re
import streamlit as st
from docx import Document
from pptx import Presentation
import pdfplumber

def extract_text(file_path):
    """通用文本提取函数"""
    try:
        if file_path.lower().endswith('.pdf'):
            with pdfplumber.open(file_path) as pdf:
                return "\n".join([p.extract_text() for p in pdf.pages if p.extract_text()])
        elif file_path.lower().endswith('.docx'):
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs if p.text])
        elif file_path.lower().endswith('.pptx'):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception as e:
        st.error(f"提取失败: {str(e)}")
        return ""

def main():
    st.title("📄 文件文字提取工具")
    file = st.file_uploader("上传文件", type=['pdf', 'docx', 'pptx', 'txt'])
    
    if file:
        with st.spinner("正在提取文字..."):
            text = extract_text(file)
        
        if text:
            st.success("提取完成！")
            st.download_button(
                "下载文本",
                text,
                file.name.split('.')[0] + ".txt"
            )
            with st.expander("预览内容"):
                st.text(text[:2000] + ("..." if len(text)>2000 else ""))

if __name__ == "__main__":
    main()