import os
import re
import pandas as pd
import streamlit as st
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import textract

# 设置页面
st.set_page_config(
    page_title="文件文字提取工具",
    page_icon="📄",
    layout="wide"
)

def extract_text_from_file(file_path):
    """根据文件类型提取文字内容"""
    try:
        if file_path.lower().endswith('.pdf'):
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                text = "\n".join([page.extract_text() for page in reader.pages])
        elif file_path.lower().endswith('.docx'):
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_path.lower().endswith('.pptx'):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            text = "\n".join(text)
        else:  # 其他格式使用textract
            text = textract.process(file_path).decode('utf-8')
        return text.strip()
    except Exception as e:
        st.error(f"提取失败: {str(e)}")
        return ""

def scan_directory(directory, extensions=['.pdf', '.docx', '.pptx', '.txt']):
    """扫描目录获取文件列表"""
    file_list = []
    for root, _, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in extensions):
                file_list.append(os.path.join(root, file))
    return sorted(file_list)

def main():
    st.title("📄 文件文字提取工具")
    st.write("从PDF/Word/PPT等文件中提取文字内容")

    # 设置扫描目录
    col1, col2 = st.columns([1, 3])
    with col1:
        st.subheader("配置选项")
        scan_dir = st.text_input(
            "要扫描的目录路径",
            value="~/360a",
            help="例如: ~/documents 或 /home/user/files"
        )
        file_extensions = st.multiselect(
            "要处理的文件类型",
            options=['PDF', 'Word', 'PPT', 'TXT', '其他'],
            default=['PDF', 'Word', 'PPT']
        )
        ext_mapping = {
            'PDF': '.pdf',
            'Word': '.docx',
            'PPT': '.pptx',
            'TXT': '.txt'
        }
        selected_exts = [ext_mapping[ext] for ext in file_extensions if ext in ext_mapping]
        
        if st.button("开始扫描目录", type="primary"):
            scan_dir = os.path.expanduser(scan_dir)
            if not os.path.exists(scan_dir):
                st.error(f"❌ 目录不存在: {scan_dir}")
                return
            
            with st.spinner("正在扫描文件..."):
                file_list = scan_directory(scan_dir, selected_exts)
            
            if not file_list:
                st.warning("⚠️ 没有找到符合条件的文件")
                return
            
     