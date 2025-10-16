import os
import tempfile
import streamlit as st
from docx import Document
from pptx import Presentation
import pdfplumber

def save_uploaded_file(uploaded_file):
    """保存上传的文件到临时文件"""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            return tmp_file.name
    except Exception as e:
        st.error(f"文件保存失败: {str(e)}")
        return None

def extract_text(file_path):
    """通用文本提取函数"""
    try:
        file_path_lower = file_path.lower()
        if file_path_lower.endswith('.pdf'):
            with pdfplumber.open(file_path) as pdf:
                return "\n".join([p.extract_text() for p in pdf.pages if p.extract_text()])
        elif file_path_lower.endswith('.docx'):
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs if p.text])
        elif file_path_lower.endswith('.pptx'):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif file_path_lower.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        else:
            st.warning(f"不支持的文件类型: {os.path.splitext(file_path)[1]}")
            return ""
    except Exception as e:
        st.error(f"提取失败: {str(e)}")
        return ""

def get_file_list(directory, extensions=['.pdf', '.docx', '.pptx', '.txt']):
    """获取目录中指定扩展名的文件列表"""
    file_list = []
    if os.path.exists(directory):
        for root, _, files in os.walk(directory):
            for file in files:
                if any(file.lower().endswith(ext) for ext in extensions):
                    file_list.append(os.path.join(root, file))
    return sorted(file_list)

def main():
    st.title("📄 文件文字提取工具")
    
    # 创建选项卡
    tab1, tab2 = st.tabs(["上传文件", "指定路径"])
    
    with tab1:
        uploaded_file = st.file_uploader("上传文件", type=['pdf', 'docx', 'pptx', 'txt'], key="uploader")
        
        if uploaded_file:
            # 保存上传的文件到临时位置
            temp_file_path = save_uploaded_file(uploaded_file)
            
            if temp_file_path:
                with st.spinner("正在提取文字..."):
                    text = extract_text(temp_file_path)
                
                # 清理临时文件
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
                
                if text:
                    st.success("提取完成！")
                    file_name = uploaded_file.name.split('.')[0] + ".txt"
                    st.download_button(
                        "下载文本",
                        text,
                        file_name,
                        key="upload_download"
                    )
                    with st.expander("预览内容", expanded=True):
                        st.text(text[:2000] + ("..." if len(text)>2000 else ""))
    
    with tab2:
        default_path = os.path.expanduser("~/360a")
        input_path = st.text_input("输入文件路径或目录", value=default_path)
        
        if st.button("扫描文件", key="scan_files"):
            if os.path.isfile(input_path):
                st.session_state.file_list = [input_path]
                st.session_state.selected_file = input_path
            elif os.path.isdir(input_path):
                st.session_state.file_list = get_file_list(input_path)
                st.session_state.selected_file = None
            else:
                st.error("路径不存在，请输入有效的文件或目录路径")
        
        if 'file_list' in st.session_state and st.session_state.file_list:
            selected_file = st.selectbox(
                "选择要提取的文件",
                options=st.session_state.file_list,
                index=0,
                key="file_selector"
            )
            
            if st.button("提取文本", key="extract_path"):
                with st.spinner("正在提取文字..."):
                    text = extract_text(selected_file)
                
                if text:
                    st.success("提取完成！")
                    file_name = os.path.basename(selected_file).split('.')[0] + ".txt"
                    st.download_button(
                        "下载文本",
                        text,
                        file_name,
                        key="path_download"
                    )
                    with st.expander("预览内容", expanded=True):
                        st.text(text[:2000] + ("..." if len(text)>2000 else ""))
        elif 'file_list' in st.session_state:
            st.warning("未找到符合条件的文件")

if __name__ == "__main__":
    main()