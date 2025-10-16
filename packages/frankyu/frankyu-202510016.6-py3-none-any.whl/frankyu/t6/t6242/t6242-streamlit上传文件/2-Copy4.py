import streamlit as st
import tempfile
import os
from PIL import Image  # 用于处理图片
import pandas as pd  # 用于处理 Excel

# 新增的库
import PyPDF2  # 用于处理 PDF 文件
from docx import Document  # 用于处理 Word (docx) 文件
from pptx import Presentation  # 用于处理 PowerPoint (pptx) 文件

# --- Streamlit 应用界面 ---
st.title("文件上传与处理应用")
st.write("请上传图片、音乐、视频、Excel、PDF、Word、PowerPoint 或文本文件，我会尝试处理并显示其内容。")

aaa = st.file_uploader("请选择一个文件进行上传", type=[
    "jpg", "jpeg", "png", "gif",            # 图片
    "mp3", "wav", "ogg",                    # 音乐
    "mp4", "mov", "avi",                    # 视频
    "xls", "xlsx",                          # Excel
    "pdf",                                  # PDF
    "doc", "docx",                          # Word (注意：只完全支持.docx)
    "ppt", "pptx",                          # PowerPoint (注意：只完全支持.pptx)
    "txt"                                   # 文本文件
])

if aaa:
    # 获取文件类型（MIME type）
    file_type = aaa.type
    st.info(f"检测到的文件类型: **{file_type}**")

    # --- 将上传文件保存到临时文件 ---
    # Streamlit 的 file_uploader 返回的文件对象通常在内存中，
    # 而许多处理文件的库需要实际的文件路径。
    # tempfile 模块允许安全地创建临时文件。
    
    # 获取原始文件的扩展名，用于临时文件
    # 注意：为了兼容性，对于Word和PPT，即使上传的是旧格式，也统一保存扩展名。
    file_extension = os.path.splitext(aaa.name)[1].lower() 
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(aaa.read())
        temp_file_path = temp_file.name

    st.success(f"文件已临时保存到: `{temp_file_path}`")

    # --- 根据文件类型选择处理方法 ---
    # 图片文件处理
    if file_type.startswith("image/"):
        st.subheader("🖼️ 图片文件处理")
        try:
            image = Image.open(temp_file_path)
            # 使用 use_container_width 替代已弃用的 use_column_width
            st.image(image, caption=f"上传的图片: {aaa.name}", use_container_width=True) 
            st.write(f"图片尺寸: **{image.size[0]} x {image.size[1]} 像素**")
            st.write(f"图片格式: **{image.format}**")
        except Exception as e:
            st.error(f"处理图片时发生错误: {e}")

    # 音乐文件处理
    elif file_type.startswith("audio/"):
        st.subheader("🎵 音乐文件处理")
        st.audio(temp_file_path, format=file_type)
        st.info("Streamlit 内置播放器将尝试播放该音频。对于更复杂的音频处理，可能需要专门的 Python 库。")

    # 视频文件处理
    elif file_type.startswith("video/"):
        st.subheader("🎬 视频文件处理")
        st.video(temp_file_path, format=file_type)
        st.info("Streamlit 内置播放器将尝试播放该视频。对于视频分析或编辑，可能需要像 OpenCV 这样的库。")

    # Excel 文件处理
    elif file_type == "application/vnd.ms-excel" or \
         file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        st.subheader("📊 Excel 文件处理")
        try:
            df = pd.read_excel(temp_file_path)
            st.success("Excel 文件读取成功！")
            st.write(f"Excel 文件包含 **{df.shape[0]} 行** 和 **{df.shape[1]} 列**。")

            st.write("文件内容预览 (前5行):")
            st.dataframe(df.head()) # 显示前5行

            # 使用 st.expander 提供查看全部内容的功能
            with st.expander("点击查看全部内容"):
                st.dataframe(df) # 显示全部内容

        except Exception as e:
            st.error(f"处理 Excel 文件时发生错误: {e}")

    # PDF 文件处理
    elif file_type == "application/pdf":
        st.subheader("📄 PDF 文件处理")
        try:
            with open(temp_file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                num_pages = len(reader.pages)
                st.write(f"PDF 文件总页数: **{num_pages}**")
                
                # 尝试提取前几页的文本内容
                extracted_text = []
                for i in range(min(num_pages, 3)): # 最多提取前3页
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        # 限制显示长度，避免过长的文本溢出
                        extracted_text.append(f"--- 第 {i+1} 页 ---\n{text[:500]}..." if len(text) > 500 else f"--- 第 {i+1} 页 ---\n{text}") 
                    else:
                        extracted_text.append(f"--- 第 {i+1} 页 ---\n(无法提取文本，可能为扫描件或图像PDF)")
                
                if extracted_text:
                    st.text_area("部分文本内容预览:", "\n\n".join(extracted_text), height=300)
                else:
                    st.info("无法从 PDF 中提取文本内容。")
                
        except Exception as e:
            st.error(f"处理 PDF 文件时发生错误: {e}")
            st.info("请注意，`PyPDF2` 主要用于文本 PDF，对于扫描版 PDF 可能无法提取文本。")

    # Word 文件处理 (仅支持 .docx)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document": # .docx MIME type
        st.subheader("📝 Word 文件处理 (.docx)")
        try:
            document = Document(temp_file_path)
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            
            doc_text = "\n".join(full_text)
            if doc_text:
                st.write(f"文档总段落数: **{len(document.paragraphs)}**")
                # 限制显示长度，避免过长的文本溢出
                st.text_area("文档内容预览 (部分):", doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text, height=300)
            else:
                st.info("无法从 Word (.docx) 文件中提取文本内容。")
        except Exception as e:
            st.error(f"处理 Word (.docx) 文件时发生错误: {e}")
            st.warning("请注意，**只完全支持 .docx 格式的 Word 文件**，不支持旧的 .doc 格式。")

    # PowerPoint 文件处理 (仅支持 .pptx)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation": # .pptx MIME type
        st.subheader("💡 PowerPoint 文件处理 (.pptx)")
        try:
            prs = Presentation(temp_file_path)
            total_slides = len(prs.slides)
            st.write(f"演示文稿总页数: **{total_slides}**")
            
            presentation_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    presentation_text.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(slide_text))

            if presentation_text:
                # 限制显示前3页的内容
                st.text_area("部分演示文稿内容预览:", "\n\n".join(presentation_text[:min(total_slides, 3)]), height=300)
            else:
                st.info("无法从 PowerPoint (.pptx) 文件中提取文本内容。")

        except Exception as e:
            st.error(f"处理 PowerPoint (.pptx) 文件时发生错误: {e}")
            st.warning("请注意，**只完全支持 .pptx 格式的 PowerPoint 文件**，不支持旧的 .ppt 格式。")

    # 文本文件处理
    elif file_type.startswith("text/"): # 涵盖 text/plain, text/csv 等
        st.subheader("📜 文本文件处理")
        try:
            with open(temp_file_path, "r", encoding="utf-8") as f:
                content = f.read()
                st.write(f"文件大小: **{len(content.encode('utf-8')) / 1024:.2f} KB**")
                # 限制显示长度，避免过长的文本溢出
                st.text_area("文件内容预览:", content[:2000] + "..." if len(content) > 2000 else content, height=300)
        except UnicodeDecodeError:
            st.warning("尝试使用 `latin-1` 编码读取文本文件...")
            try:
                with open(temp_file_path, "r", encoding="latin-1") as f:
                    content = f.read()
                    st.text_area("文件内容预览:", content[:2000] + "..." if len(content) > 2000 else content, height=300)
            except Exception as e:
                st.error(f"处理文本文件时发生编码错误: {e}")
        except Exception as e:
            st.error(f"处理文本文件时发生错误: {e}")

    # 其他未支持的文件类型
    else:
        st.warning(f"🤔 抱歉，当前不支持处理类型为 `{file_type}` 的文件。")

    # --- 清理临时文件 ---
    # 确保文件处理完成后删除临时文件，避免占用磁盘空间。
    os.unlink(temp_file_path)
    st.info("临时文件已删除。")