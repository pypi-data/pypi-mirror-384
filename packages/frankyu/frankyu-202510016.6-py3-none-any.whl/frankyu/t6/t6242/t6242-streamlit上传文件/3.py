#pip install streamlit pandas PyPDF2 python-docx python-pptx openpyxl xlrd textract pypandoc pillow   -i   https://pypi.tuna.tsinghua.edu.cn/simple



# 忽略特定警告
import warnings
warnings.filterwarnings("ignore", category=UserWarning)

# --- 尝试导入各函数库，并追踪其可用性 ---
libraries_available = {
    "pandas": False,
    "PyPDF2": False,
    "docx": False,
    "pptx": False,
    "openpyxl": False,
    "xlrd": False,
    "textract": False,
    "pypandoc": False
}

import streamlit as st
import tempfile
import os
import shutil
import time
import io
import subprocess

# 尝试导入 pandas
try:
    import pandas as pd
    libraries_available["pandas"] = True
except ImportError:
    pass

# 尝试导入 PyPDF2
try:
    import PyPDF2
    libraries_available["PyPDF2"] = True
except ImportError:
    pass

# 尝试导入 python-docx
try:
    from docx import Document
    libraries_available["docx"] = True
except ImportError:
    pass

# 尝试导入 python-pptx
try:
    from pptx import Presentation
    libraries_available["pptx"] = True
except ImportError:
    pass

# 尝试导入 openpyxl
try:
    import openpyxl
    libraries_available["openpyxl"] = True
except ImportError:
    pass

# 尝试导入 xlrd (用于旧版 Excel)
try:
    import xlrd
    libraries_available["xlrd"] = True
except ImportError:
    pass

# 尝试导入 textract (用于提取文本)
try:
    import textract
    libraries_available["textract"] = True
except ImportError:
    pass

# 尝试导入 pypandoc (用于文档转换)
try:
    import pypandoc
    # 确保 pandoc 已安装
    try:
        pypandoc.get_pandoc_path()
        libraries_available["pypandoc"] = True
    except OSError:
        pass
except ImportError:
    pass

# Pillow 函数库检查 - 必需库
try:
    from PIL import Image
except ImportError:
    st.error("❌ **严重错误：** 缺少 `Pillow` 函数库。\n"
             "图片显示与处理功能将无法使用。\n"
             "请运行 `pip install Pillow`。")
    st.stop()

# --- 显示所有缺失库的警告 ---
if not libraries_available["pandas"]:
    st.warning("⚠️ **警告：** 缺少 `pandas` 函数库。\n"
               "Excel 文件 (`.xls`, `.xlsx`) 处理功能将无法使用。\n"
               "请运行 `pip install pandas`。")

if not libraries_available["PyPDF2"]:
    st.warning("⚠️ **警告：** 缺少 `PyPDF2` 函数库。\n"
               "PDF 文件处理功能将无法使用。\n"
               "请运行 `pip install PyPDF2`。")

if not libraries_available["docx"]:
    st.warning("⚠️ **警告：** 缺少 `python-docx` 函数库。\n"
               "Word (.docx) 文件处理功能将无法使用。\n"
               "请运行 `pip install python-docx`。")

if not libraries_available["pptx"]:
    st.warning("⚠️ **警告：** 缺少 `python-pptx` 函数库。\n"
               "PowerPoint (.pptx) 文件处理功能将无法使用。\n"
               "请运行 `pip install python-pptx`。")

if not libraries_available["openpyxl"]:
    if libraries_available["pandas"]:
        st.warning("⚠️ **警告：** 缺少 `openpyxl` 函数库。\n"
                   "Excel 的高级功能（如图片提取）可能受限。\n"
                   "请运行 `pip install openpyxl`。")

if not libraries_available["xlrd"]:
    st.warning("⚠️ **警告：** 缺少 `xlrd` 函数库。\n"
               "旧版 Excel (.xls) 文件处理功能将无法使用。\n"
               "请运行 `pip install xlrd`。")

if not libraries_available["textract"]:
    st.warning("⚠️ **警告：** 缺少 `textract` 函数库。\n"
               "部分文件格式的文本提取功能将受限。\n"
               "请运行 `pip install textract`。")

if not libraries_available["pypandoc"]:
    st.warning("⚠️ **警告：** 缺少 `pypandoc` 函数库或 `pandoc` 命令行工具。\n"
               "文档转换功能将受限。\n"
               "请安装:\n"
               "1. `pip install pypandoc`\n"
               "2. 从 https://pandoc.org/installing.html 安装 pandoc")

# --- Streamlit 应用程序界面 ---
st.set_page_config(layout="wide")
st.title("通用文件上传与处理工具")
st.write("请上传图片、音讯、视讯、Excel、PDF、Word、PowerPoint 或文字文件，\n"
         "我会尝试处理并显示其内容。")

uploaded_file = st.file_uploader("选择一个文件上传", type=[
    "jpg", "jpeg", "png", "gif",            # 图片
    "mp3", "wav", "ogg",                    # 音讯
    "mp4", "mov", "avi",                    # 视讯
    "xls", "xlsx",                          # Excel
    "pdf",                                  # PDF
    "doc", "docx",                          # Word
    "ppt", "pptx",                          # PowerPoint
    "txt"                                   # 文字文件
])

# 初始化变量
temp_file_path = None
converted_file_path = None

def convert_with_pandoc(input_path, output_extension):
    """使用 pandoc 转换文件格式"""
    try:
        output_path = tempfile.mktemp(suffix=output_extension)
        pypandoc.convert_file(
            input_path,
            output_extension[1:],  # 格式如 'docx'
            outputfile=output_path
        )
        return output_path
    except Exception as e:
        st.error(f"文件转换失败: {str(e)}")
        return None

def extract_text_with_textract(file_path):
    """使用 textract 提取文件文本"""
    try:
        text = textract.process(file_path).decode('utf-8')
        return text
    except Exception as e:
        st.error(f"文本提取失败: {str(e)}")
        return None

def is_pandoc_available():
    """检查 pandoc 是否可用"""
    try:
        result = subprocess.run(["pandoc", "--version"], capture_output=True, text=True)
        return result.returncode == 0
    except:
        return False

if uploaded_file:
    # 获取文件的 MIME 类型
    file_type = uploaded_file.type
    st.info(f"检测到的文件类型: **{file_type}**")
    st.info(f"文件名称: **{uploaded_file.name}**")

    # 保存上传文件到临时位置
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name

    st.success(f"文件已暂存至: `{temp_file_path}`")

    # 根据文件类型处理文件
    # 图片文件
    if file_type.startswith("image/"):
        st.subheader("🖼️ 图片文件处理")
        try:
            image = Image.open(temp_file_path)
            st.image(image,
                     caption=f"上传图片: {uploaded_file.name}",
                     use_container_width=True)
            st.write(f"图片尺寸: **{image.size[0]} x {image.size[1]} 像素**")
            st.write(f"图片格式: **{image.format}**")
            image.close()
        except Exception as e:
            st.error(f"处理图片文件时发生错误: {e}")

    # 音讯文件
    elif file_type.startswith("audio/"):
        st.subheader("🎵 音讯文件处理")
        st.audio(temp_file_path, format=file_type)
        st.info("Streamlit 内置播放器会尝试播放音讯。\n"
                "若需更复杂的音讯处理，需使用专用 Python 函数库。")

    # 视讯文件
    elif file_type.startswith("video/"):
        st.subheader("🎬 视讯文件处理")
        st.video(temp_file_path, format=file_type)
        st.info("Streamlit 内置播放器会尝试播放视讯。\n"
                "若需视讯分析或编辑，需使用如 OpenCV 等函数库。")

    # Excel 文件 (xls, xlsx)
    elif file_type == "application/vnd.ms-excel" or \
         file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        st.subheader("📊 Excel 文件处理")
        
        # 处理旧版 .xls 文件
        if file_extension == ".xls":
            if libraries_available["xlrd"]:
                try:
                    # 使用 xlrd 读取旧版 Excel
                    workbook = xlrd.open_workbook(temp_file_path)
                    sheet_names = workbook.sheet_names()
                    
                    st.success("Excel 文件读取成功！")
                    st.write(f"此 Excel 文件包含 **{len(sheet_names)} 个工作表**。")
                    
                    selected_sheet = st.selectbox("选择要查看的工作表:",
                                                  sheet_names,
                                                  index=0,
                                                  key="excel_sheet_select")
                    
                    sheet = workbook.sheet_by_name(selected_sheet)
                    
                    # 显示表格数据
                    data = []
                    for row_idx in range(0, min(sheet.nrows, 100)):
                        row_data = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            row_data.append(cell_value)
                        data.append(row_data)
                    
                    st.write(f"**显示工作表: `{selected_sheet}` 的内容 (预览前 100 行)**")
                    st.dataframe(data, use_container_width=True)
                    
                    # 提取文本内容
                    text_content = []
                    for row in data:
                        for cell in row:
                            if isinstance(cell, (str, int, float)) and str(cell).strip():
                                text_content.append(str(cell).strip())
                    
                    if text_content:
                        with st.expander("点击查看所有提取的单元格文字"):
                            st.text_area("提取的文字:",
                                         "\n---\n".join(text_content),
                                         height=300)
                    else:
                        st.info("此工作表的单元格中未找到任何显著的文字内容。")
                    
                except Exception as e:
                    st.error(f"处理旧版 Excel 文件时发生错误: {e}")
            else:
                st.warning("缺少 `xlrd` 库，无法处理旧版 Excel (.xls) 文件。")
        
        # 处理新版 .xlsx 文件
        elif file_extension == ".xlsx" and libraries_available["pandas"]:
            excel_file_handle = None
            temp_image_dir = None
            try:
                excel_file_handle = pd.ExcelFile(temp_file_path)
                sheet_names = excel_file_handle.sheet_names

                st.success("Excel 文件读取成功！")
                st.write(f"此 Excel 文件包含 **{len(sheet_names)} 个工作表**。")

                selected_sheet = st.selectbox("选择要查看的工作表:",
                                              sheet_names,
                                              index=0,
                                              key="excel_sheet_select")

                df_selected = pd.read_excel(excel_file_handle,
                                            sheet_name=selected_sheet)

                st.write(f"**显示工作表: `{selected_sheet}` 的内容 (预览前 100 行)**")
                st.dataframe(df_selected.head(100), use_container_width=True)

                if len(df_selected) > 100:
                    if st.checkbox(f"点击查看 `{selected_sheet}` 的完整表格\n"
                                   "(可能需要更多内存和时间)",
                                   key=f"show_full_df_{selected_sheet}"):
                        st.dataframe(df_selected, use_container_width=True)

                st.subheader(f"📝 工作表: `{selected_sheet}` 中的文字内容")
                text_content = []
                for col in df_selected.columns:
                    for cell_value in df_selected[col].astype(str):
                        if pd.notna(cell_value) and len(cell_value.strip()) > 0:
                            text_content.append(cell_value.strip())

                if text_content:
                    with st.expander("点击查看所有提取的单元格文字"):
                        st.text_area("提取的文字:",
                                     "\n---\n".join(text_content),
                                     height=300)
                else:
                    st.info("此工作表的单元格中未找到任何显著的文字内容。")

                # 图片提取部分
                if libraries_available["openpyxl"]:
                    st.subheader(f"🖼️ 工作表: `{selected_sheet}` 中的图片")
                    image_found = False
                    try:
                        temp_image_dir = tempfile.mkdtemp()
                        workbook = openpyxl.load_workbook(temp_file_path)
                        
                        # 使用更健壮的方式获取工作表
                        try:
                            sheet = workbook[selected_sheet]
                        except KeyError:
                            sheet = None
                            for ws in workbook.worksheets:
                                if ws.title == selected_sheet:
                                    sheet = ws
                                    break
                            
                            if sheet is None:
                                st.error(f"找不到名称为 '{selected_sheet}' 的工作表")
                                images = []
                            else:
                                images = getattr(sheet, '_images', [])
                        else:
                            images = getattr(sheet, '_images', [])
                        
                        if images:
                            total_images = len(images)
                            max_images_to_preview = 5

                            st.write(f"在 `{selected_sheet}` 中找到 **{total_images}** 张图片。")

                            if total_images > max_images_to_preview:
                                st.info(f"为了性能，**仅预览前 {max_images_to_preview} 张图片**。")
                                if st.checkbox("显示所有图片 (可能导致卡顿，请谨慎)",
                                               key=f"show_all_images_{selected_sheet}"):
                                    max_images_to_preview = total_images

                            for i, img in enumerate(images):
                                if i >= max_images_to_preview:
                                    break

                                # 尝试确定图片扩展名
                                image_ext = ".png"
                                if hasattr(img, 'format') and img.format:
                                    if img.format.lower() in ['png', 'jpeg', 'gif']:
                                        image_ext = f".{img.format.lower()}"
                                
                                image_filename = f"image_{i+1}{image_ext}"
                                image_path = os.path.join(temp_image_dir, image_filename)

                                try:
                                    # 正确获取图片数据
                                    if hasattr(img, '_data'):
                                        pil_image = Image.open(io.BytesIO(img._data()))
                                        
                                        # 转换图片模式以确保兼容性
                                        if pil_image.mode in ('RGBA', 'P'):
                                            pil_image = pil_image.convert('RGB')
                                        
                                        # 保存图片
                                        pil_image.save(image_path)
                                        st.image(image_path,
                                                 caption=f"图片 {i+1} 来自 {selected_sheet}",
                                                 use_container_width=True)
                                        image_found = True
                                    else:
                                        st.warning(f"图片 {i+1} 缺少数据，无法显示")
                                except Exception as img_err:
                                    st.warning(f"无法保存图片 {i+1}。错误: {img_err}")

                            if image_found:
                                st.success(f"成功提取并显示 {min(total_images, max_images_to_preview)} 张图片")
                            else:
                                st.info("此工作表中找到图片但无法显示")
                        else:
                            st.info("此工作表中未找到任何嵌入图片。")

                    except Exception as e:
                        st.error(f"从Excel文件中提取图片时发生错误: {e}")
                    finally:
                        if temp_image_dir and os.path.exists(temp_image_dir):
                            try:
                                shutil.rmtree(temp_image_dir)
                            except Exception as e:
                                st.warning(f"无法删除临时图片目录 {temp_image_dir}。原因: {e}")
                else:
                    st.warning("`openpyxl`库未加载，无法进行Excel图片提取。")
            except Exception as e:
                st.error(f"处理 Excel 文件时发生错误: {e}")
            finally:
                if excel_file_handle:
                    excel_file_handle.close()
        else:
            st.warning("无法处理此 Excel 文件格式。")

    # PDF 文件
    elif file_type == "application/pdf":
        st.subheader("📄 PDF 文件处理")
        if libraries_available["PyPDF2"]:
            pdf_file = None
            try:
                pdf_file = open(temp_file_path, "rb")
                reader = PyPDF2.PdfReader(pdf_file)
                num_pages = len(reader.pages)
                st.write(f"PDF 文件总页数: **{num_pages}**")

                extracted_text = []
                for i in range(min(num_pages, 3)):
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        extracted_text.append(f"--- 第 {i+1} 页 ---\n"
                                              f"{text[:500]}..." if len(text) > 500 else
                                              f"--- 第 {i+1} 页 ---\n{text}")
                    else:
                        extracted_text.append(f"--- 第 {i+1} 页 ---\n"
                                              "(无法提取文字，可能为扫描件或图像式 PDF)")

                if extracted_text:
                    st.text_area("部分文字内容预览:",
                                 "\n\n".join(extracted_text),
                                 height=300)
                else:
                    st.info("无法从 PDF 中提取任何文字内容。")

            except Exception as e:
                st.error(f"处理 PDF 文件时发生错误: {e}")
                st.info("注意: `PyPDF2` 主要用于文字型 PDF，\n"
                        "可能无法从扫描或图像型 PDF 中提取文字。")
            finally:
                if pdf_file:
                    pdf_file.close()
        else:
            st.warning("`PyPDF2` 函数库未加载，无法处理 PDF 文件。")

    # Word 文件 (.doc 和 .docx)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or \
         file_type == "application/msword":
        st.subheader("📝 Word 文件处理")
        
        # 处理 .docx 文件
        if file_extension == ".docx" and libraries_available["docx"]:
            try:
                document = Document(temp_file_path)
                full_text = []
                for para in document.paragraphs:
                    full_text.append(para.text)

                doc_text = "\n".join(full_text)
                if doc_text:
                    st.write(f"文件总段落数: **{len(document.paragraphs)}**")
                    st.text_area("文件内容预览:",
                                 doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text,
                                 height=300)
                else:
                    st.info("无法从 Word 文件中提取任何文字内容。")
            except Exception as e:
                st.error(f"处理 Word 文件时发生错误: {e}")
        
        # 处理旧版 .doc 文件
        elif file_extension == ".doc":
            st.warning("检测到旧版 Word (.doc) 文件，尝试提取文本内容...")
            
            # 尝试使用 pandoc 转换为 docx
            if libraries_available["pypandoc"] and is_pandoc_available():
                st.info("尝试使用 pandoc 转换文件...")
                try:
                    converted_file_path = convert_with_pandoc(temp_file_path, ".docx")
                    if converted_file_path and os.path.exists(converted_file_path):
                        st.success("文件转换成功！")
                        
                        # 尝试处理转换后的 docx 文件
                        try:
                            document = Document(converted_file_path)
                            full_text = []
                            for para in document.paragraphs:
                                full_text.append(para.text)

                            doc_text = "\n".join(full_text)
                            if doc_text:
                                st.text_area("文件内容预览:",
                                             doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text,
                                             height=300)
                            else:
                                st.info("无法从转换后的文件中提取任何文字内容。")
                        except Exception as e:
                            st.error(f"处理转换后的文件时发生错误: {e}")
                except Exception as e:
                    st.error(f"文件转换失败: {str(e)}")
            
            # 如果 pandoc 不可用，尝试使用 textract
            elif libraries_available["textract"]:
                st.info("尝试使用 textract 提取文本...")
                try:
                    text = extract_text_with_textract(temp_file_path)
                    if text:
                        st.text_area("提取的文本内容:",
                                     text[:2000] + "..." if len(text) > 2000 else text,
                                     height=300)
                    else:
                        st.info("无法从 .doc 文件中提取文本内容。")
                except Exception as e:
                    st.error(f"提取文本时出错: {e}")
            
            else:
                st.warning("无法处理旧版 .doc 文件。请安装:\n"
                           "1. `pip install textract` (需要系统安装 antiword)\n"
                           "2. `pip install pypandoc` 并安装 pandoc")
        
        else:
            st.warning("无法处理此 Word 文件格式。")

    # PowerPoint 文件 (.ppt 和 .pptx)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or \
         file_type == "application/vnd.ms-powerpoint":
        st.subheader("💡 PowerPoint 文件处理")
        
        # 处理 .pptx 文件
        if file_extension == ".pptx" and libraries_available["pptx"]:
            try:
                prs = Presentation(temp_file_path)
                total_slides = len(prs.slides)
                st.write(f"简报总页数: **{total_slides}**")

                presentation_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        presentation_text.append(f"--- 投影片 {i+1} ---\n" + "\n".join(slide_text))

                if presentation_text:
                    st.text_area("部分简报内容预览:",
                                 "\n\n".join(presentation_text[:min(total_slides, 3)]),
                                 height=300)
                else:
                    st.info("无法从 PowerPoint 文件中提取任何文字内容。")
            except Exception as e:
                st.error(f"处理 PowerPoint 文件时发生错误: {e}")
        
        # 处理旧版 .ppt 文件
        elif file_extension == ".ppt":
            st.warning("检测到旧版 PowerPoint (.ppt) 文件，尝试提取文本内容...")
            
            # 尝试使用 textract 提取文本
            if libraries_available["textract"]:
                st.info("尝试使用 textract 提取文本...")
                try:
                    text = extract_text_with_textract(temp_file_path)
                    if text:
                        st.text_area("提取的文本内容:",
                                     text[:2000] + "..." if len(text) > 2000 else text,
                                     height=300)
                    else:
                        st.info("无法从 .ppt 文件中提取文本内容。")
                except Exception as e:
                    st.error(f"提取文本时出错: {e}")
            else:
                st.warning("无法处理旧版 .ppt 文件。请安装 textract 库：`pip install textract`")
        
        else:
            st.warning("无法处理此 PowerPoint 文件格式。")

    # 文字文件
    elif file_type.startswith("text/"):
        st.subheader("📜 文字文件处理")
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'latin-1', 'cp1252', 'gbk', 'gb2312', 'big5']
            content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(temp_file_path, "r", encoding=encoding) as text_file:
                        content = text_file.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                st.error("无法使用任何常见编码读取文件")
            else:
                st.success(f"使用 {used_encoding} 编码成功读取文件")
                st.write(f"文件大小: **{len(content.encode('utf-8')) / 1024:.2f} KB**")
                st.text_area("文件内容预览:",
                             content[:2000] + "..." if len(content) > 2000 else content,
                             height=300)
        except Exception as e:
            st.error(f"处理文字文件时发生错误: {e}")

    # 其他不支持的档案类型
    else:
        st.warning(f"🤔 抱歉，目前不支持处理类型为 `{file_type}` 的档案。")

    # --- 清理暂存文件 ---
    if temp_file_path and os.path.exists(temp_file_path):
        try:
            os.unlink(temp_file_path)
        except Exception as e:
            st.warning(f"删除原始暂存文件时出错: {e}")
    
    if converted_file_path and os.path.exists(converted_file_path):
        try:
            os.unlink(converted_file_path)
        except Exception as e:
            st.warning(f"删除转换后文件时出错: {e}")

else:
    st.info("👆 请上传文件以开始处理")