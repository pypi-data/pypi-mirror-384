import streamlit as st
import tempfile
import os
from PIL import Image # Pillow 函式庫，用於圖片處理
import shutil # 用於刪除非空目錄

# --- 嘗試導入各函式庫，並追蹤其可用性 ---
# 建立一個字典來追蹤每個函式庫的載入狀態
libraries_available = {
    "pandas": False,
    "PyPDF2": False,
    "docx": False, # 對應 python-docx
    "pptx": False, # 對應 python-pptx
    "openpyxl": False
}

# 嘗試導入 pandas
try:
    import pandas as pd
    libraries_available["pandas"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `pandas` 函式庫。\n"
               "Excel 檔案 (`.xls`, `.xlsx`) 處理功能將無法使用。\n"
               "請運行 `pip install pandas`。")

# 嘗試導入 PyPDF2
try:
    import PyPDF2
    libraries_available["PyPDF2"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `PyPDF2` 函式庫。\n"
               "PDF 檔案處理功能將無法使用。\n"
               "請運行 `pip install PyPDF2`。")

# 嘗試導入 python-docx
try:
    from docx import Document
    libraries_available["docx"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `python-docx` 函式庫。\n"
               "Word (.docx) 檔案處理功能將無法使用。\n"
               "請運行 `pip install python-docx`。")

# 嘗試導入 python-pptx
try:
    from pptx import Presentation
    libraries_available["pptx"] = True
except ImportError:
    st.warning("⚠️ **警告：** 缺少 `python-pptx` 函式庫。\n"
               "PowerPoint (.pptx) 檔案處理功能將無法使用。\n"
               "請運行 `pip install python-pptx`。")

# 嘗試導入 openpyxl
# openpyxl 通常與 pandas 一起安裝，但為確保健壯性，仍單獨檢查
# 只有在 pandas 載入成功的前提下，openpyxl 的導入才有意義 (因為用於 Excel 圖片)
if libraries_available["pandas"]:
    try:
        import openpyxl
        libraries_available["openpyxl"] = True
    except ImportError:
        st.warning("⚠️ **警告：** 缺少 `openpyxl` 函式庫。\n"
                   "Excel 的高級功能（如圖片提取）可能受限。\n"
                   "請運行 `pip install openpyxl`。")
else:
    # 如果 pandas 都沒裝，那麼 openpyxl 即使裝了也無法完全發揮作用，直接設為不可用
    libraries_available["openpyxl"] = False

# Pillow 函式庫 (PIL) 是基礎圖片處理功能，如果它都無法載入，則直接停止應用程式
try:
    # 這裡的 Image 已經在頂部導入，如果上面沒有報錯，說明 PIL 已經可用
    # 這裡只是作為一個最終確認，如果確實無法載入，則停止
    pass
except ImportError:
    st.error("❌ **嚴重錯誤：** 缺少 `Pillow` 函式庫。\n"
             "圖片顯示與處理功能將無法使用。\n"
             "請運行 `pip install Pillow`。")
    st.stop() # 如果最核心的圖片庫都缺失，則無法繼續運行

# --- Streamlit 應用程式介面 ---
st.set_page_config(layout="wide") # 設置頁面佈局為寬模式，更好顯示表格
st.title("通用文件上傳與處理工具")
st.write("請上傳圖片、音訊、視訊、Excel、PDF、Word、PowerPoint 或文字檔案，\n"
         "我會嘗試處理並顯示其內容。")

uploaded_file = st.file_uploader("選擇一個檔案上傳", type=[
    "jpg", "jpeg", "png", "gif",            # 圖片
    "mp3", "wav", "ogg",                    # 音訊
    "mp4", "mov", "avi",                    # 視訊
    "xls", "xlsx",                          # Excel
    "pdf",                                  # PDF
    "doc", "docx",                          # Word (注意：僅完全支援 .docx)
    "ppt", "pptx",                          # PowerPoint (注意：僅完全支援 .pptx)
    "txt"                                   # 文字檔案
])

# 初始化 temp_file_path，確保它總是有定義，以便最後進行清理
temp_file_path = None

if uploaded_file:
    # 取得檔案的 MIME 類型
    file_type = uploaded_file.type
    st.info(f"檢測到的檔案類型: **{file_type}**")

    # --- 將上傳檔案儲存到暫存位置 ---
    # Streamlit 的 file_uploader 通常將檔案儲存在記憶體中，
    # 但許多處理函式庫需要實際的檔案路徑。
    # tempfile 模組可以安全地建立暫存檔案。

    # 取得原始檔案的副檔名，用於暫存檔案
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()

    # 使用 with 語句確保暫存檔案在寫入完成後被正確關閉
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name # 在這裡指定路徑

    st.success(f"檔案已暫存至: `{temp_file_path}`")

    # --- 根據檔案類型處理檔案 ---

    # 圖片檔案
    if file_type.startswith("image/"):
        st.subheader("🖼️ 圖片檔案處理")
        # Pillow 已經在頂部強制檢查，如果應用程式運行到這裡，Pillow 應該可用
        try:
            image = Image.open(temp_file_path)
            # 使用 use_container_width 以適應容器寬度
            st.image(image,
                     caption=f"上傳圖片: {uploaded_file.name}",
                     use_container_width=True)
            st.write(f"圖片尺寸: **{image.size[0]} x {image.size[1]} 像素**")
            st.write(f"圖片格式: **{image.format}**")
            image.close() # 明確關閉圖片檔案句柄
        except Exception as e:
            st.error(f"處理圖片檔案時發生錯誤: {e}")

    # 音訊檔案
    elif file_type.startswith("audio/"):
        st.subheader("🎵 音訊檔案處理")
        st.audio(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放音訊。\n"
                "若需更複雜的音訊處理，需使用專用 Python 函式庫。")

    # 視訊檔案
    elif file_type.startswith("video/"):
        st.subheader("🎬 視訊檔案處理")
        st.video(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放視訊。\n"
                "若需視訊分析或編輯，需使用如 OpenCV 等函式庫。")

    # Excel 檔案 (xls, xlsx) - 增強多工作表、文字與圖片處理
    elif file_type == "application/vnd.ms-excel" or \
         file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        st.subheader("📊 Excel 檔案處理")
        if libraries_available["pandas"]: # 檢查 pandas 是否可用
            excel_file_handle = None
            temp_image_dir = None # 初始化臨時圖片目錄變數
            try:
                excel_file_handle = pd.ExcelFile(temp_file_path)
                sheet_names = excel_file_handle.sheet_names

                st.success("Excel 檔案讀取成功！")
                st.write(f"此 Excel 檔案包含 **{len(sheet_names)} 個工作表**。")

                # --- 工作表選擇 ---
                selected_sheet = st.selectbox("選擇要查看的工作表:",
                                              sheet_names,
                                              index=0,
                                              key="excel_sheet_select")

                df_selected = pd.read_excel(excel_file_handle,
                                            sheet_name=selected_sheet)

                # --- 優化：限制表格預覽行數 ---
                st.write(f"**顯示工作表: `{selected_sheet}` 的內容 (預覽前 100 行)**")
                st.dataframe(df_selected.head(100), use_container_width=True)

                if len(df_selected) > 100:
                    if st.checkbox(f"點擊查看 `{selected_sheet}` 的完整表格\n"
                                   "(可能需要更多記憶體和時間)",
                                   key=f"show_full_df_{selected_sheet}"):
                        st.dataframe(df_selected, use_container_width=True)

                # --- 從選定工作表提取文字內容 ---
                st.subheader(f"📝 工作表: `{selected_sheet}` 中的文字內容")
                text_content = []
                for col in df_selected.columns:
                    # 將所有值轉換為字串，以便迭代和檢查是否為空
                    for cell_value in df_selected[col].astype(str):
                        if pd.notna(cell_value) and len(cell_value.strip()) > 0:
                            text_content.append(cell_value.strip())

                if text_content:
                    with st.expander("點擊查看所有提取的儲存格文字"):
                        st.text_area("提取的文字:",
                                     "\n---\n".join(text_content),
                                     height=300)
                else:
                    st.info("此工作表的儲存格中未找到任何顯著的文字內容。")

                # --- 圖片提取 (僅限 .xlsx 檔案) ---
                if file_extension == ".xlsx":
                    st.subheader(f"🖼️ 工作表: `{selected_sheet}` 中的圖片")
                    if libraries_available["openpyxl"]: # 檢查 openpyxl 是否可用
                        image_found = False
                        try:
                            # 建立一個臨時目錄用於儲存圖片
                            temp_image_dir = tempfile.mkdtemp()

                            workbook = openpyxl.load_workbook(temp_file_path)
                            sheet = workbook[selected_sheet]

                            # openpyxl 圖片存儲在 _images 屬性中
                            # 注意: _images 是一個內部屬性，但在 openpyxl 中是提取圖片的常用方式。
                            if hasattr(sheet, '_images') and sheet._images: # 檢查是否存在 _images 且非空
                                total_images = len(sheet._images)
                                max_images_to_preview = 5 # 預設只預覽前5張圖片

                                st.write(f"在 `{selected_sheet}` 中找到 **{total_images}** 張圖片。")

                                if total_images > max_images_to_preview:
                                    st.info(f"為了效能，**僅預覽前 {max_images_to_preview} 張圖片**。")
                                    if st.checkbox("顯示所有圖片 (可能導致卡頓，請謹慎)",
                                                   key=f"show_all_images_{selected_sheet}"):
                                        max_images_to_preview = total_images # 允許顯示所有圖片

                                for i, img in enumerate(sheet._images):
                                    if i >= max_images_to_preview:
                                        break # 超出預覽限制，停止預覽

                                    image_ext = ".png" # 預設副檔名

                                    # 更穩健地獲取圖片副檔名
                                    if hasattr(img, 'ext') and img.ext:
                                        image_ext = f".{img.ext}"
                                    elif hasattr(img.ref, 'embed'):
                                        # 嘗試從 img.ref.embed (ImagePart) 中獲取圖片的 MIME 類型
                                        # 並根據 MIME 類型判斷副檔名
                                        img_data = img.ref.embed
                                        if img_data and hasattr(img_data, 'mime_type') and img_data.mime_type:
                                            mime_type = img_data.mime_type
                                            if 'image/png' in mime_type:
                                                image_ext = '.png'
                                            elif 'image/jpeg' in mime_type or 'image/jpg' in mime_type:
                                                image_ext = '.jpg'
                                            elif 'image/gif' in mime_type:
                                                image_ext = '.gif'
                                            elif 'image/bmp' in mime_type:
                                                image_ext = '.bmp'
                                            elif 'image/tiff' in mime_type:
                                                image_ext = '.tiff'
                                            # 您可以根據需要添加更多 mime 類型

                                    image_filename = f"image_{i}{image_ext}"
                                    image_path = os.path.join(temp_image_dir, image_filename)

                                    try:
                                        # 將 openpyxl 提供的圖片數據流傳遞給 PIL.Image.open
                                        pil_image = Image.open(img.ref)
                                        # 優化：儲存為 JPEG 格式，並降低品質以減少檔案大小
                                        if pil_image.mode in ('RGBA', 'P'): # JPEG 不支援透明度，通常轉換為 RGB
                                            pil_image = pil_image.convert('RGB')
                                        pil_image.save(image_path, format='JPEG', quality=85) # quality 0-100

                                        st.image(image_path,
                                                 caption=f"圖片 {i+1} 來自 {selected_sheet}",
                                                 use_container_width=True)
                                        image_found = True
                                    except Exception as img_err:
                                        st.warning(f"無法儲存圖片 {i+1} 到 {image_path}。\n"
                                                   f"錯誤: {img_err}\n"
                                                   "這可能是因為圖片格式不被 PIL 支援，或者圖片數據有問題。")

                            if not image_found:
                                st.info("此工作表中未找到任何嵌入圖片。")

                        except Exception as e:
                            st.error(f"從 Excel 檔案中提取圖片時發生錯誤\n"
                                     f"(圖片提取僅支援 .xlsx): {e}")
                        finally:
                            # 清理臨時圖片目錄
                            if temp_image_dir and os.path.exists(temp_image_dir):
                                try:
                                    shutil.rmtree(temp_image_dir)
                                except Exception as e:
                                    st.warning(f"無法刪除臨時圖片目錄 {temp_image_dir}。\n"
                                               f"原因: {e}")
                    else:
                        st.warning("`openpyxl` 函式庫未載入，無法進行 Excel 圖片提取。\n"
                                   "請檢查之前的警告訊息以了解如何安裝。")
                else:
                    st.info("圖片提取僅支援 .xlsx 檔案。")

            except Exception as e:
                st.error(f"處理 Excel 檔案時發生錯誤: {e}")
            finally:
                # 確保 ExcelFile 句柄被明確關閉
                if excel_file_handle:
                    excel_file_handle.close()
        else:
            st.warning("`pandas` 函式庫未載入，無法處理 Excel 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # PDF 檔案
    elif file_type == "application/pdf":
        st.subheader("📄 PDF 檔案處理")
        if libraries_available["PyPDF2"]: # 檢查 PyPDF2 是否可用
            pdf_file = None
            try:
                pdf_file = open(temp_file_path, "rb")
                reader = PyPDF2.PdfReader(pdf_file)
                num_pages = len(reader.pages)
                st.write(f"PDF 檔案總頁數: **{num_pages}**")

                extracted_text = []
                for i in range(min(num_pages, 3)): # 最多提取前 3 頁的文字
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        # 限制顯示長度以避免文字溢出
                        extracted_text.append(f"--- 第 {i+1} 頁 ---\n"
                                              f"{text[:500]}..." if len(text) > 500 else
                                              f"--- 第 {i+1} 頁 ---\n{text}")
                    else:
                        extracted_text.append(f"--- 第 {i+1} 頁 ---\n"
                                              "(無法提取文字，可能為掃描件或圖像式 PDF)")

                if extracted_text:
                    st.text_area("部分文字內容預覽:",
                                 "\n\n".join(extracted_text),
                                 height=300)
                else:
                    st.info("無法從 PDF 中提取任何文字內容。")

            except Exception as e:
                st.error(f"處理 PDF 檔案時發生錯誤: {e}")
                st.info("注意: `PyPDF2` 主要用於文字型 PDF，\n"
                        "可能無法從掃描或圖像型 PDF 中提取文字。")
            finally:
                if pdf_file:
                    pdf_file.close()
        else:
            st.warning("`PyPDF2` 函式庫未載入，無法處理 PDF 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # Word 檔案 (僅限 DOCX)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        st.subheader("📝 Word 檔案處理 (.docx)")
        if libraries_available["docx"]: # 檢查 docx 是否可用
            try:
                document = Document(temp_file_path)
                full_text = []
                for para in document.paragraphs:
                    full_text.append(para.text)

                doc_text = "\n".join(full_text)
                if doc_text:
                    st.write(f"文件總段落數: **{len(document.paragraphs)}**")
                    st.text_area("文件內容預覽 (部分):",
                                 doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text,
                                 height=300)
                else:
                    st.info("無法從 Word (.docx) 檔案中提取任何文字內容。")
            except Exception as e:
                st.error(f"處理 Word (.docx) 檔案時發生錯誤: {e}")
                st.warning("請注意: **僅完全支援 .docx 格式的 Word 檔案**，\n"
                           "不支援舊的 .doc 格式。")
        else:
            st.warning("`python-docx` 函式庫未載入，無法處理 Word 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # PowerPoint 檔案 (僅限 PPTX)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        st.subheader("💡 PowerPoint 檔案處理 (.pptx)")
        if libraries_available["pptx"]: # 檢查 pptx 是否可用
            try:
                prs = Presentation(temp_file_path)
                total_slides = len(prs.slides)
                st.write(f"簡報總頁數: **{total_slides}**")

                presentation_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        presentation_text.append(f"--- 投影片 {i+1} ---\n" + "\n".join(slide_text))

                if presentation_text:
                    st.text_area("部分簡報內容預覽:",
                                 "\n\n".join(presentation_text[:min(total_slides, 3)]),
                                 height=300)
                else:
                    st.info("無法從 PowerPoint (.pptx) 檔案中提取任何文字內容。")
            except Exception as e:
                st.error(f"處理 PowerPoint (.pptx) 檔案時發生錯誤: {e}")
                st.warning("請注意: **僅完全支援 .pptx 格式的 PowerPoint 檔案**，\n"
                           "不支援舊的 .ppt 格式。")
        else:
            st.warning("`python-pptx` 函式庫未載入，無法處理 PowerPoint 檔案。\n"
                       "請檢查之前的警告訊息以了解如何安裝。")

    # 文字檔案
    elif file_type.startswith("text/"):
        st.subheader("📜 文字檔案處理")
        text_file = None
        try:
            text_file = open(temp_file_path, "r", encoding="utf-8")
            content = text_file.read()
            st.write(f"檔案大小: **{len(content.encode('utf-8')) / 1024:.2f} KB**")
            st.text_area("檔案內容預覽:",
                         content[:2000] + "..." if len(content) > 2000 else content,
                         height=300)
        except UnicodeDecodeError:
            st.warning("嘗試使用 `latin-1` 編碼讀取文字檔案...")
            try:
                text_file = open(temp_file_path, "r", encoding="latin-1")
                content = text_file.read()
                st.text_area("檔案內容預覽:",
                             content[:2000] + "..." if len(content) > 2000 else content,
                             height=300)
            except Exception as e:
                st.error(f"處理文字檔案時發生編碼錯誤: {e}")
        except Exception as e:
            st.error(f"處理文字檔案時發生錯誤: {e}")
        finally:
            if text_file:
                text_file.close()

    # 其他不支援的檔案類型
    else:
        st.warning(f"🤔 抱歉，目前不支援處理類型為 `{file_type}` 的檔案。")

    # --- 清理暫存檔案 ---
    if temp_file_path and os.path.exists(temp_file_path):
        try:
            os.unlink(temp_file_path)
            st.info("暫存檔案已成功刪除。")
        except PermissionError:
            st.warning("⚠️ **權限錯誤：** 無法刪除暫存檔案。\n"
                       "檔案可能仍被系統或其他程式使用中。\n"
                       "請稍後手動刪除或如果問題持續，請重新啟動應用程式。")
        except Exception as e:
            st.error(f"刪除暫存檔案時發生未知錯誤: {e}")
else:
    pass # 沒有上傳檔案時不執行任何操作