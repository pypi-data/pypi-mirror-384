import streamlit as st
import tempfile
import os
from PIL import Image  # 用於圖片處理
import pandas as pd  # 用於 Excel 處理

# 額外的文件處理函式庫
import PyPDF2  # 用於 PDF 檔案
from docx import Document  # 用於 Word (.docx) 檔案
from pptx import Presentation  # 用於 PowerPoint (.pptx) 檔案

# --- Streamlit 應用程式介面 ---
st.title("通用文件上傳與處理工具")
st.write("請上傳圖片、音訊、視訊、Excel、PDF、Word、PowerPoint 或文字檔案。我會嘗試處理並顯示其內容。")

uploaded_file = st.file_uploader("選擇一個檔案上傳", type=[
    "jpg", "jpeg", "png", "gif",            # 圖片
    "mp3", "wav", "ogg",                    # 音訊
    "mp4", "mov", "avi",                    # 視訊
    "xls", "xlsx",                          # Excel
    "pdf",                                  # PDF
    "doc", "docx",                          # Word (注意：僅完全支援 .docx)
    "ppt", "pptx",                          # PowerPoint (注意：僅完全支援 .pptx)
    "txt"                                   # 文字檔案
])

# 初始化 temp_file_path，確保在任何情況下都能定義它，以便最後進行清理
temp_file_path = None 

if uploaded_file:
    # 取得檔案的 MIME 類型
    file_type = uploaded_file.type
    st.info(f"檢測到的檔案類型: **{file_type}**")

    # --- 將上傳檔案儲存到暫存位置 ---
    # Streamlit 的 file_uploader 通常將檔案儲存在記憶體中，
    # 但許多處理函式庫需要實際的檔案路徑。
    # tempfile 模組可以安全地建立暫存檔案。
    
    # 取得原始檔案的副檔名，用於暫存檔案
    file_extension = os.path.splitext(uploaded_file.name)[1].lower() 
    
    # 使用 with 語句確保暫存檔案在寫入完成後被正確關閉
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name # 在這裡指定路徑

    st.success(f"檔案已暫存至: `{temp_file_path}`")

    # --- 根據檔案類型處理檔案 ---

    # 圖片檔案
    if file_type.startswith("image/"):
        st.subheader("🖼️ 圖片檔案處理")
        try:
            image = Image.open(temp_file_path)
            # 使用 use_container_width 替代已棄用的 use_column_width
            st.image(image, caption=f"上傳圖片: {uploaded_file.name}", use_container_width=True) 
            st.write(f"圖片尺寸: **{image.size[0]} x {image.size[1]} 像素**")
            st.write(f"圖片格式: **{image.format}**")
            image.close() # 明確關閉圖片檔案句柄
        except Exception as e:
            st.error(f"處理圖片檔案時發生錯誤: {e}")

    # 音訊檔案
    elif file_type.startswith("audio/"):
        st.subheader("🎵 音訊檔案處理")
        st.audio(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放音訊。若需更複雜的音訊處理，需使用專用 Python 函式庫。")

    # 視訊檔案
    elif file_type.startswith("video/"):
        st.subheader("🎬 視訊檔案處理")
        st.video(temp_file_path, format=file_type)
        st.info("Streamlit 內建播放器會嘗試播放視訊。若需視訊分析或編輯，需使用如 OpenCV 等函式庫。")

    # Excel 檔案 (xls, xlsx) - 增強多工作表與穩健處理
    elif file_type == "application/vnd.ms-excel" or \
         file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        st.subheader("📊 Excel 檔案處理")
        excel_file_handle = None # 初始化為 None，以便管理檔案關閉
        try:
            # 使用 pd.ExcelFile 妥善管理 Excel 檔案，並確保之後關閉它
            excel_file_handle = pd.ExcelFile(temp_file_path)
            sheet_names = excel_file_handle.sheet_names
            
            st.success("Excel 檔案讀取成功！")
            st.write(f"此 Excel 檔案包含 **{len(sheet_names)} 個工作表**。")
            
            # 預設顯示第一個工作表的頭部內容
            default_sheet_name = sheet_names[0]
            st.write(f"**目前顯示工作表: `{default_sheet_name}` (前 5 列)**")
            # 將 ExcelFile 物件傳遞給 pd.read_excel
            df_default = pd.read_excel(excel_file_handle, sheet_name=default_sheet_name) 
            st.dataframe(df_default.head()) 

            # 如果存在多個工作表，允許使用者選擇其他工作表
            if len(sheet_names) > 1: 
                selected_sheet = st.selectbox("選擇要查看完整內容的工作表:", sheet_names, index=0) # 預設選中第一個
                
                if selected_sheet:
                    st.write(f"**顯示工作表: `{selected_sheet}` 的完整內容**")
                    # 將 ExcelFile 物件傳遞給 pd.read_excel
                    df_selected = pd.read_excel(excel_file_handle, sheet_name=selected_sheet) 
                    st.dataframe(df_selected, use_container_width=True) 
            else:
                # 如果只有一個工作表，提供一個可展開區塊來查看完整內容
                with st.expander("點擊查看完整內容"):
                    st.dataframe(df_default, use_container_width=True) 

        except Exception as e:
            st.error(f"處理 Excel 檔案時發生錯誤: {e}")
        finally:
            # 確保 ExcelFile 句柄被明確關閉
            if excel_file_handle:
                excel_file_handle.close()
                # st.write("Excel 檔案句柄已關閉。") # 可用於調試

    # PDF 檔案
    elif file_type == "application/pdf":
        st.subheader("📄 PDF 檔案處理")
        pdf_file = None # 初始化為 None，以便明確關閉
        try:
            # 明確打開檔案供 PyPDF2 讀取
            pdf_file = open(temp_file_path, "rb") 
            reader = PyPDF2.PdfReader(pdf_file)
            num_pages = len(reader.pages)
            st.write(f"PDF 檔案總頁數: **{num_pages}**")
            
            extracted_text = []
            for i in range(min(num_pages, 3)): # 最多提取前 3 頁的文字
                page = reader.pages[i]
                text = page.extract_text()
                if text:
                    # 限制顯示長度以避免文字溢出
                    extracted_text.append(f"--- 第 {i+1} 頁 ---\n{text[:500]}..." if len(text) > 500 else f"--- 第 {i+1} 頁 ---\n{text}") 
                else:
                    extracted_text.append(f"--- 第 {i+1} 頁 ---\n(無法提取文字，可能為掃描件或圖像式 PDF)")
                
            if extracted_text:
                st.text_area("部分文字內容預覽:", "\n\n".join(extracted_text), height=300)
            else:
                st.info("無法從 PDF 中提取任何文字內容。")
                
        except Exception as e:
            st.error(f"處理 PDF 檔案時發生錯誤: {e}")
            st.info("注意: `PyPDF2` 主要用於文字型 PDF，可能無法從掃描或圖像型 PDF 中提取文字。")
        finally:
            # 確保 PDF 檔案句柄已關閉
            if pdf_file:
                pdf_file.close() 

    # Word 檔案 (僅限 DOCX)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document": # .docx 的 MIME 類型
        st.subheader("📝 Word 檔案處理 (.docx)")
        try:
            document = Document(temp_file_path)
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            
            doc_text = "\n".join(full_text)
            if doc_text:
                st.write(f"文件總段落數: **{len(document.paragraphs)}**")
                # 限制顯示長度以避免文字溢出
                st.text_area("文件內容預覽 (部分):", doc_text[:1000] + "..." if len(doc_text) > 1000 else doc_text, height=300)
            else:
                st.info("無法從 Word (.docx) 檔案中提取任何文字內容。")
            # python-docx 函式庫通常會自動處理檔案關閉，或不會持有持久性鎖定。
        except Exception as e:
            st.error(f"處理 Word (.docx) 檔案時發生錯誤: {e}")
            st.warning("請注意: **僅完全支援 .docx 格式的 Word 檔案**，不支援舊的 .doc 格式。")

    # PowerPoint 檔案 (僅限 PPTX)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation": # .pptx 的 MIME 類型
        st.subheader("💡 PowerPoint 檔案處理 (.pptx)")
        try:
            prs = Presentation(temp_file_path)
            total_slides = len(prs.slides)
            st.write(f"簡報總頁數: **{total_slides}**")
            
            presentation_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    presentation_text.append(f"--- 投影片 {i+1} ---\n" + "\n".join(slide_text))

            if presentation_text:
                # 限制顯示前 3 頁的內容
                st.text_area("部分簡報內容預覽:", "\n\n".join(presentation_text[:min(total_slides, 3)]), height=300)
            else:
                st.info("無法從 PowerPoint (.pptx) 檔案中提取任何文字內容。")
            # python-pptx 函式庫通常會自動處理檔案關閉，或不會持有持久性鎖定。
        except Exception as e:
            st.error(f"處理 PowerPoint (.pptx) 檔案時發生錯誤: {e}")
            st.warning("請注意: **僅完全支援 .pptx 格式的 PowerPoint 檔案**，不支援舊的 .ppt 格式。")

    # 文字檔案
    elif file_type.startswith("text/"): # 包含 text/plain, text/csv 等
        st.subheader("📜 文字檔案處理")
        text_file = None # 初始化以便明確關閉
        try:
            # 明確打開檔案進行讀取
            text_file = open(temp_file_path, "r", encoding="utf-8") 
            content = text_file.read()
            st.write(f"檔案大小: **{len(content.encode('utf-8')) / 1024:.2f} KB**")
            # 限制顯示長度以避免文字溢出
            st.text_area("檔案內容預覽:", content[:2000] + "..." if len(content) > 2000 else content, height=300)
        except UnicodeDecodeError:
            st.warning("嘗試使用 `latin-1` 編碼讀取文字檔案...")
            try:
                # 如果 UTF-8 失敗，則再次嘗試以 latin-1 編碼打開
                text_file = open(temp_file_path, "r", encoding="latin-1") 
                content = text_file.read()
                st.text_area("檔案內容預覽:", content[:2000] + "..." if len(content) > 2000 else content, height=300)
            except Exception as e:
                st.error(f"處理文字檔案時發生編碼錯誤: {e}")
        except Exception as e:
            st.error(f"處理文字檔案時發生錯誤: {e}")
        finally:
            # 確保文字檔案句柄已關閉
            if text_file:
                text_file.close() 

    # 其他不支援的檔案類型
    else:
        st.warning(f"🤔 抱歉，目前不支援處理類型為 `{file_type}` 的檔案。")

    # --- 清理暫存檔案 ---
    # 此區塊在所有檔案處理嘗試後執行。
    # 特定的檔案句柄現在已在其各自的區塊中明確關閉。
    if temp_file_path and os.path.exists(temp_file_path):
        try:
            os.unlink(temp_file_path)
            st.info("暫存檔案已成功刪除。")
        except PermissionError:
            st.warning("⚠️ 無法刪除暫存檔案。檔案可能仍被系統或其他程式使用中。請稍後手動刪除或如果問題持續，請重新啟動應用程式。")
        except Exception as e:
            st.error(f"刪除暫存檔案時發生未知錯誤: {e}")
else:
    # 如果沒有檔案上傳，此區塊確保 temp_file_path 不用於刪除邏輯。
    pass