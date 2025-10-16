# ragchunker/file_ingestion.py
"""
File Ingestion Layer for ragchunker.

Provides:
- Document dataclass (content + metadata)
- Per-file loaders: PDF, DOCX, PPTX, TXT/MD, Images (OCR), Audio (Whisper), Excel/CSV/Parquet, HTML, JSON, ZIP/TAR
- auto_load(path) function to detect file type and return list[Document]

Notes:
- Many features are optional and rely on external libraries. The loader will raise informative errors
  or return an empty list if a dependency is absent.
- Audio transcription uses OpenAI's whisper package (local model) by default if installed.
"""

from __future__ import annotations
import os
import io
import mimetypes
import zipfile
import tarfile
import uuid
import json
import csv
import logging
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Any, Union

# Third-party imports with friendly fallbacks in loader functions
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except Exception:
    pdf_extract_text = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

try:
    import whisper  # openai-whisper
except Exception:
    whisper = None

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from bs4 import BeautifulSoup
except Exception:
    BeautifulSoup = None

# Optional: trafilatura for robust HTML extraction
try:
    import trafilatura
except Exception:
    trafilatura = None

# configure logger for the module
logger = logging.getLogger("ragchunker.file_ingestion")
if not logger.handlers:
    # Basic config if not already configured by the main app
    logging.basicConfig(level=logging.INFO)


@dataclass
class Document:
    """
    Represents an extracted document chunk or unit of text with metadata.
    """
    id: str = field(default_factory=lambda: str(uuid.uuid4()))
    content: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)


# --------------------------
# Helper utilities
# --------------------------
def _read_text_file(path: str, encoding: str = "utf-8") -> str:
    with open(path, "r", encoding=encoding, errors="ignore") as f:
        return f.read()


def _flatten_dataframe_to_text(df: "pd.DataFrame") -> List[Document]:
    """
    Convert a DataFrame to a list of documents (one per row).
    Each row will be stringified as 'col1: val1\ncol2: val2' etc.
    """
    docs: List[Document] = []
    for idx, row in df.iterrows():
        parts = []
        for col in df.columns:
            parts.append(f"{col}: {row[col]}")
        text = "\n".join(parts)
        docs.append(Document(content=text, metadata={"row_index": int(idx)}))
    return docs


def _safe_mimetype(path: str) -> str:
    mt, _ = mimetypes.guess_type(path)
    return mt or "application/octet-stream"


# --------------------------
# Loaders
# --------------------------
def load_txt(path: str) -> List[Document]:
    try:
        text = _read_text_file(path)
        return [Document(content=text, metadata={"source": path, "type": "text"})]
    except Exception as e:
        logger.exception("Failed to load text file %s: %s", path, e)
        return []


def load_markdown(path: str) -> List[Document]:
    # treat like text for now; could add markdown-aware parsing later
    return load_txt(path)


def load_pdf(path: str, page_level: bool = True) -> List[Document]:
    """
    Extract text from PDF. If page_level True, returns one Document per page (if supported),
    otherwise returns a single Document with whole text.
    """
    if pdf_extract_text is None:
        raise ImportError("pdfminer.six is required for PDF extraction. Install: pip install pdfminer.six")

    try:
        # pdfminer extract_text returns the whole text; it doesn't provide page separation by default.
        # A simple approach: extract whole and return single doc. If page-level needed, user can supply library like pdfplumber.
        whole_text = pdf_extract_text(path)
        if not whole_text:
            logger.warning("PDF no text extracted for %s", path)
            return []
        return [Document(content=whole_text, metadata={"source": path, "type": "pdf"})]
    except Exception as e:
        logger.exception("Failed to extract PDF %s: %s", path, e)
        return []


def load_docx(path: str) -> List[Document]:
    if docx is None:
        raise ImportError("python-docx is required for DOCX. Install: pip install python-docx")
    try:
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        content = "\n".join(paragraphs)
        return [Document(content=content, metadata={"source": path, "type": "docx"})]
    except Exception as e:
        logger.exception("Failed to read DOCX %s: %s", path, e)
        return []


def load_pptx(path: str) -> List[Document]:
    if Presentation is None:
        raise ImportError("python-pptx is required for PPTX. Install: pip install python-pptx")
    docs: List[Document] = []
    try:
        prs = Presentation(path)
        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text
                if text and text.strip():
                    texts.append(text.strip())
            slide_text = "\n".join(texts)
            docs.append(Document(content=slide_text, metadata={"source": path, "type": "pptx", "slide_index": slide_idx}))
        return docs
    except Exception as e:
        logger.exception("Failed to read PPTX %s: %s", path, e)
        return []


def load_image(path: str, ocr_lang: Optional[str] = None) -> List[Document]:
    """
    Extract text from image using pytesseract. Returns a single Document containing OCR result.
    """
    if Image is None or pytesseract is None:
        raise ImportError("Pillow and pytesseract required for image OCR. Install: pip install pillow pytesseract")
    try:
        img = Image.open(path)
        # optionally convert to RGB/grayscale to improve OCR reliability
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        ocr_config = None
        text = pytesseract.image_to_string(img, lang=ocr_lang, config=ocr_config) if ocr_lang else pytesseract.image_to_string(img)
        if not text:
            logger.warning("No OCR text detected in %s", path)
        return [Document(content=text, metadata={"source": path, "type": "image", "ocr_lang": ocr_lang})]
    except Exception as e:
        logger.exception("Failed to OCR image %s: %s", path, e)
        return []


def load_audio_whisper(path: str, model_name: str = "small") -> List[Document]:
    """
    Transcribe audio using local Whisper (openai-whisper package).
    Model names: tiny, base, small, medium, large.
    Returns a list of Document objects (optionally one per segment).
    """

    import shutil

    if whisper is None:
        raise ImportError(
            "The 'whisper' package is required for audio transcription. "
            "Install with: pip install -U openai-whisper"
        )

    # Check file exists
    if not os.path.isfile(path):
        logger.error("Audio file not found: %s", path)
        return []

    # Check ffmpeg availability
    if shutil.which("ffmpeg") is None:
        logger.error(
            "ffmpeg not found in PATH. Whisper depends on ffmpeg to load audio files.\n"
            "Install ffmpeg (conda, winget, choco, or manual) and restart your terminal.\n"
            "Examples:\n"
            "  conda install -c conda-forge ffmpeg\n"
            "  winget install --id=Gyan.FFmpeg -e\n"
            "  choco install ffmpeg -y\n"
            "Or download a release and add the ffmpeg 'bin' folder to your PATH."
        )
        return []

    try:
        # This will download the model if not cached — be aware of disk/network usage
        model = whisper.load_model(model_name)
        result = model.transcribe(path)

        text = result.get("text", "")
        segments = result.get("segments", None)
        metadata = {"source": path, "type": "audio", "whisper_model": model_name}

        if segments:
            docs = []
            for seg in segments:
                seg_text = seg.get("text", "")
                docs.append(
                    Document(
                        content=seg_text,
                        metadata={
                            **metadata,
                            "start": seg.get("start"),
                            "end": seg.get("end"),
                        },
                    )
                )
            return docs

        return [Document(content=text, metadata=metadata)]

    except FileNotFoundError as e:
        # Likely ffmpeg subprocess not found — but we've already tested for ffmpeg above.
        logger.exception("FileNotFoundError during transcription of %s: %s", path, e)
        return []
    except Exception as e:
        logger.exception("Failed to transcribe audio %s: %s", path, e)
        return []



def load_csv(path: str, encoding: str = "utf-8") -> List[Document]:
    if pd is None:
        raise ImportError("pandas required for CSV/Excel/Parquet loaders. Install: pip install pandas pyarrow")
    try:
        df = pd.read_csv(path, encoding=encoding, dtype=str)
        docs = _flatten_dataframe_to_text(df)
        for d in docs:
            d.metadata.update({"source": path, "type": "csv"})
        return docs
    except Exception as e:
        logger.exception("Failed to read CSV %s: %s", path, e)
        return []


def load_excel(path: str, sheet_name: Union[str, int, None] = None) -> List[Document]:
    if pd is None:
        raise ImportError("pandas required for CSV/Excel/Parquet loaders. Install: pip install pandas openpyxl pyarrow")
    try:
        df = pd.read_excel(path, sheet_name=sheet_name, dtype=str)
        # pandas returns dict when sheet_name=None; convert to single list
        if isinstance(df, dict):
            docs: List[Document] = []
            for sheet, subdf in df.items():
                part_docs = _flatten_dataframe_to_text(subdf)
                for d in part_docs:
                    d.metadata.update({"source": path, "type": "excel", "sheet": str(sheet)})
                docs.extend(part_docs)
            return docs
        else:
            docs = _flatten_dataframe_to_text(df)
            for d in docs:
                d.metadata.update({"source": path, "type": "excel"})
            return docs
    except Exception as e:
        logger.exception("Failed to read Excel %s: %s", path, e)
        return []


def load_parquet(path: str) -> List[Document]:
    if pd is None:
        raise ImportError("pandas and pyarrow are required for Parquet. Install: pip install pandas pyarrow")
    try:
        df = pd.read_parquet(path)
        docs = _flatten_dataframe_to_text(df)
        for d in docs:
            d.metadata.update({"source": path, "type": "parquet"})
        return docs
    except Exception as e:
        logger.exception("Failed to read Parquet %s: %s", path, e)
        return []


def load_html(path: str) -> List[Document]:
    """
    Load and extract visible text from HTML. Prefer trafilatura if available; else fallback to BeautifulSoup.
    """
    try:
        with open(path, "rb") as f:
            raw = f.read()
    except Exception as e:
        logger.exception("Failed to read HTML file %s: %s", path, e)
        return []

    if trafilatura:
        try:
            text = trafilatura.extract(raw.decode("utf-8", errors="ignore"))
            if text:
                return [Document(content=text, metadata={"source": path, "type": "html", "engine": "trafilatura"})]
        except Exception:
            logger.warning("trafilatura failed on %s", path)

    if BeautifulSoup:
        try:
            soup = BeautifulSoup(raw, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator="\n")
            return [Document(content=text, metadata={"source": path, "type": "html", "engine": "bs4"})]
        except Exception as e:
            logger.exception("BeautifulSoup failed on %s: %s", path, e)
            return []

    raise ImportError("Either trafilatura or beautifulsoup4 is required for HTML extraction. Install: pip install trafilatura beautifulsoup4")


def load_json(path: str) -> List[Document]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        logger.exception("Failed to load JSON %s: %s", path, e)
        return []

    # Flatten JSON to text. For lists, create doc per item; for dict, create one doc per top-level key
    docs: List[Document] = []
    if isinstance(data, list):
        for idx, item in enumerate(data):
            docs.append(Document(content=json.dumps(item, ensure_ascii=False), metadata={"source": path, "type": "json", "index": idx}))
    elif isinstance(data, dict):
        for k, v in data.items():
            docs.append(Document(content=json.dumps({k: v}, ensure_ascii=False), metadata={"source": path, "type": "json", "key": k}))
    else:
        docs.append(Document(content=str(data), metadata={"source": path, "type": "json"}))
    return docs


def extract_from_zip_or_tar(path: str, tmp_extract_dir: Optional[str] = None) -> List[str]:
    """
    Extract archive to a temp dir (or a sibling folder) and return list of extracted file paths.
    Does not delete extracted files. Caller should manage cleanup.
    """
    extracted_paths: List[str] = []
    base_dir = tmp_extract_dir or (os.path.splitext(path)[0] + "_extracted")
    os.makedirs(base_dir, exist_ok=True)
    try:
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path, "r") as z:
                z.extractall(base_dir)
                extracted_paths = [os.path.join(base_dir, p) for p in z.namelist()]
        else:
            # try tar
            try:
                with tarfile.open(path, "r:*") as tar:
                    tar.extractall(base_dir)
                    extracted_paths = [os.path.join(base_dir, m.name) for m in tar.getmembers() if m.isreg()]
            except Exception as e:
                logger.exception("Archive extraction failed for %s: %s", path, e)
                return []
    except Exception as e:
        logger.exception("Error extracting archive %s: %s", path, e)
        return []
    # Filter out directories
    extracted_paths = [p for p in extracted_paths if os.path.isfile(p)]
    return extracted_paths


# --------------------------
# Auto-detection loader
# --------------------------
EXT_MAP = {
    ".txt": load_txt,
    ".md": load_markdown,
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".webp": load_image,
    ".bmp": load_image,
    ".gif": load_image,
    ".mp3": load_audio_whisper,
    ".wav": load_audio_whisper,
    ".m4a": load_audio_whisper,
    ".flac": load_audio_whisper,
    ".csv": load_csv,
    ".xls": load_excel,
    ".xlsx": load_excel,
    ".parquet": load_parquet,
    ".html": load_html,
    ".htm": load_html,
    ".json": load_json,
    ".zip": extract_from_zip_or_tar,
    ".tar": extract_from_zip_or_tar,
    ".tar.gz": extract_from_zip_or_tar,
    ".tgz": extract_from_zip_or_tar,
}


def auto_load(path: str, **kwargs) -> List[Document]:
    """
    Auto-detect the path and dispatch to appropriate loader.
    - If path is directory: recursively loads supported files within.
    - If path is archive: extracts and loads files within.
    - kwargs forwarded to loaders (e.g., chunkers may want ocr_lang, whisper model_name)
    """
    path = os.path.abspath(path)
    if not os.path.exists(path):
        raise FileNotFoundError(f"{path} does not exist")

    docs: List[Document] = []

    if os.path.isdir(path):
        # Walk directory
        for root, _, files in os.walk(path):
            for fname in files:
                full = os.path.join(root, fname)
                docs.extend(auto_load(full, **kwargs))
        return docs

    _, ext = os.path.splitext(path)
    ext = ext.lower()

    # special-case archive extensions with multiple suffixes
    if any(path.lower().endswith(suffix) for suffix in (".tar.gz", ".tar.bz2", ".tgz", ".tar")):
        ext = os.path.splitext(path)[1].lower()
        # fallback: treat as archive
        loader = extract_from_zip_or_tar
    else:
        loader = EXT_MAP.get(ext)

    # If no loader found by ext, use mimetype best-effort dispatch
    if loader is None:
        mt = _safe_mimetype(path)
        logger.debug("Unknown extension %s, guessed mimetype %s", ext, mt)
        if mt and mt.startswith("text"):
            loader = load_txt
        elif mt and mt.startswith("image"):
            loader = load_image
        elif mt and mt.startswith("audio"):
            loader = load_audio_whisper
        else:
            # fallback to text read
            loader = load_txt

    try:
        result = loader(path, **kwargs) if loader is not extract_from_zip_or_tar else loader(path, kwargs.get("tmp_extract_dir"))
    except TypeError:
        # if loader signature doesn't accept kwargs
        result = loader(path)

    # If the loader returned paths (archive extractor), recursively load those
    if isinstance(result, list) and result and all(isinstance(x, str) and os.path.exists(x) for x in result):
        for p in result:
            docs.extend(auto_load(p, **kwargs))
    elif isinstance(result, list) and all(isinstance(x, Document) for x in result):
        # ensure metadata contains source where not present
        for d in result:
            if "source" not in d.metadata:
                d.metadata.setdefault("source", path)
        docs.extend(result)
    elif result is None:
        # loader handled error and returned None
        logger.debug("Loader returned None for %s", path)
    else:
        # unexpected return type (e.g., single Document)
        if isinstance(result, Document):
            result.metadata.setdefault("source", path)
            docs.append(result)
        else:
            logger.debug("Loader returned unexpected type %s for %s", type(result), path)
    return docs


# --------------------------
# Convenience API
# --------------------------
def load_single(path: str, prefer_single_doc: bool = True, **kwargs) -> Document:
    """
    Load a path and return a single combined Document. If multiple docs are extracted,
    they will be concatenated (with separators).
    """
    docs = auto_load(path, **kwargs)
    if not docs:
        return Document(content="", metadata={"source": path})
    if len(docs) == 1:
        return docs[0]
    # combine
    combined_text = "\n\n---\n\n".join([d.content for d in docs if d.content])
    combined_meta = {"source": path, "parts": len(docs)}
    return Document(content=combined_text, metadata=combined_meta)


# --------------------------
# If run as script - quick demo
# --------------------------
if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser(description="Quick demo for ragchunker file ingestion")
    p.add_argument("path", help="File or directory to ingest")
    p.add_argument("--ocr_lang", default=None, help="pytesseract OCR language (optional)")
    p.add_argument("--whisper_model", default="small", help="Whisper model name (tiny, base, small, medium, large)")
    args = p.parse_args()

    print(f"Loading {args.path} ...")
    try:
        # Forward whisper model via kwargs when loading audio
        docs = auto_load(args.path, ocr_lang=args.ocr_lang, model_name=args.whisper_model)
        print(f"Extracted {len(docs)} documents")
        for i, doc in enumerate(docs[:10]):
            print(f"\n--- Document {i+1} (id={doc.id}) ---")
            print("Metadata:", doc.metadata)
            text_preview = doc.content[:1000].replace("\n", " ")
            print("Preview:", text_preview)
    except Exception as e:
        logger.exception("Failed ingestion demo: %s", e)
