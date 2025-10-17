import io

from bs4 import BeautifulSoup
from pdfminer.layout import LAParams

try:
    from pdfminer.high_level import extract_text

    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

from kagura.core.models import ModelRegistry, validate_required_state_fields

StateModel = ModelRegistry.get("StateModel")


class TextConverterError(Exception):
    pass


class TextConverter:
    _instance = None

    def __init__(self):
        pass

    def __new__(cls, *args, **kwargs):
        if not cls._instance:
            cls._instance = super(TextConverter, cls).__new__(cls)
        return cls._instance

    @classmethod
    def get_instance(cls):
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    def extract_text_content_from_pdf(self, pdf_content: bytes) -> str:
        if not HAS_PDF:
            raise TextConverterError("pdfminer.six not installed")
        pdf_io = io.BytesIO(pdf_content)
        laparams = LAParams(
            line_margin=0.5,
            word_margin=0.1,
            boxes_flow=0.5,
            detect_vertical=True,  # detect vertical text
            all_texts=True,
        )
        text = extract_text(
            pdf_io,
            laparams=laparams,
            codec="utf-8",
        )
        return text.strip()

    def extract_text_content_from_pptx(self, pptx_content: bytes) -> str:
        if not HAS_PPTX:
            raise TextConverterError("python-pptx not installed")
        pptx_io = io.BytesIO(pptx_content)
        presentation = Presentation(pptx_io)
        return "\n\n".join(
            shape.text.strip()
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )

    def extract_text_content_from_docx(self, docx_content: bytes) -> str:
        if not HAS_DOCX:
            raise TextConverterError("python-docx not installed")
        docx_io = io.BytesIO(docx_content)
        document = Document(docx_io)
        return "\n\n".join(
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )

    def extract_text_content_from_html(self, html_content: str) -> str:
        try:
            soup = BeautifulSoup(html_content, "html.parser")
            for tag in soup.find_all(["script", "style", "nav", "footer", "header"]):
                tag.decompose()

            main_content = None
            for selector in [
                "main",
                "article",
                '[role="main"]',
                "#main-content",
                "#content",
            ]:
                main_content = soup.select_one(selector)
                if main_content:
                    break

            if not main_content:
                main_content = soup.body or soup

            return "\n\n".join(
                p.get_text(strip=True)
                for p in main_content.find_all(
                    ["p", "h1", "h2", "h3", "h4", "h5", "h6"]
                )
                if p.get_text(strip=True)
            )
        except Exception as e:
            raise TextConverterError(f"HTML parsing failed: {str(e)}")

    async def convert(self, content: str, content_type: str) -> str:
        handlers = {
            "pdf": self.extract_text_content_from_pdf,
            "pptx": self.extract_text_content_from_pptx,
            "docx": self.extract_text_content_from_docx,
            "html": self.extract_text_content_from_html,
        }

        handler = handlers.get(content_type)
        if not handler:
            return content

        try:
            return handler(content)
        except Exception as e:
            raise TextConverterError(f"Conversion failed for {content_type}: {str(e)}")


async def convert_to_text(state: StateModel) -> StateModel:
    validate_required_state_fields(state, ["content", "converted_content"])

    try:
        state.converted_content = await TextConverter.get_instance().convert(
            state.content.text, state.content.content_type
        )
        return state
    except Exception as e:
        raise TextConverterError(str(e))
