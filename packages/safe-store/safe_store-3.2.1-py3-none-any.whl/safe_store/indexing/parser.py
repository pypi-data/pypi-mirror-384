# safe_store/indexing/parser.py
from pathlib import Path
from typing import Callable, Dict, Union, List
from ascii_colors import ASCIIColors
import os
import pipmaster as pm
# Import specific custom exceptions
from ..core.exceptions import ParsingError, FileHandlingError, ConfigurationError

# Protocols remain unchanged, add type hints
from typing import Protocol, runtime_checkable, BinaryIO, TextIO
@runtime_checkable
class PdfReaderProtocol(Protocol): # noqa
    pages: list
    def __init__(self, stream: Union[bytes, Path, str, BinaryIO], strict: bool = True): ...
@runtime_checkable
class PageObjectProtocol(Protocol): # noqa
    def extract_text(self) -> str: ...
@runtime_checkable
class DocumentProtocol(Protocol): # noqa
    paragraphs: list
    def __init__(self, file_path: Union[str, Path, None] = None): ...
@runtime_checkable
class ParagraphProtocol(Protocol): # noqa
    text: str
@runtime_checkable
class BeautifulSoupProtocol(Protocol): # noqa
    def __init__(self, markup: Union[str, bytes], features: str, **kwargs): ...
    def get_text(self, separator: str = "", strip: bool = False) -> str: ...


# --- Parsing Functions ---
def _process_msg_attachment(data: bytes, filename: str, images: dict) -> str:
    # Guess file type based on filename
    mime_type, _ = mimetypes.guess_type(filename)
    file_ext = filename.lower().split(".")[-1] if "." in filename else ""

    # Try decoding as text if likely to be a text attachment
    if mime_type and mime_type.startswith("text") or file_ext in ("txt", "csv", "md"):
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = data.decode("latin1")
            except Exception:
                text = None
        if text and text.strip():
            return f"[Attachment: {filename}]\n{text[:300]}{'...' if len(text)>300 else ''}"
        else:
            return f"[Attachment: {filename} -- undecodable text]"

    # If it's a PDF, Word, or other known binary type, just summarize
    if file_ext in ("pdf", "docx", "xlsx", "pptx", "zip", "eml", "msg") or \
       (mime_type and (mime_type.startswith("application/") or mime_type.startswith("image/"))):
        # Optionally store images for further processing
        if mime_type and mime_type.startswith("image") and images is not None:
            images[filename] = data  # Store image bytes for future use
        return f"[Attachment: {filename} ({mime_type or file_ext}) not parsed as text]"

    # Otherwise, unknown type
    return f"[Attachment: {filename} (unknown type) -- size {len(data)} bytes]"

def parse_msg(file_path: Union[str, Path]) -> str:
    try:
        import extract_msg
        from bs4 import BeautifulSoup  # Correct import
    except ImportError:
        pm.ensure_packages(["extract-msg", "bs4"])  # Use 'bs4' not 'BeautifulSoup'
        import extract_msg
        from bs4 import BeautifulSoup

    try:
        msg = extract_msg.Message(str(file_path))
        images = {}
        header_lines = []
        if getattr(msg, "subject", ""):
            header_lines.append(f"# {getattr(msg, 'subject', '')}")
        meta = []
        sender = getattr(msg, "sender", None) or getattr(msg, "from_", None)
        if sender:
            meta.append(f"From: {sender}")
        if getattr(msg, "to", ""):
            meta.append(f"To: {getattr(msg, 'to', '')}")
        if getattr(msg, "date", ""):
            meta.append(f"Date: {getattr(msg, 'date', '')}")
        if meta:
            header_lines.append("\n".join(meta))
        header = "\n\n".join(header_lines)

        msg_body = (getattr(msg, "body", "") or "").strip()
        if not msg_body and getattr(msg, "htmlBody", None):
            html_body = getattr(msg, "htmlBody", "")
            if html_body:
                msg_body = BeautifulSoup(html_body, "html.parser").get_text()

        attachment_text_parts: List[str] = [header, msg_body]
        if hasattr(msg, "attachments"):
            for att in msg.attachments:
                att_data = getattr(att, "data", b"")
                att_name = getattr(att, "longFilename", None) or getattr(att, "shortFilename", None) or "attachment"
                text_part = _process_msg_attachment(att_data, att_name, images)
                if text_part:
                    attachment_text_parts.append(text_part)

        extracted_text = "\n\n".join([p for p in attachment_text_parts if p and p.strip()])
        return extracted_text
    except Exception as e:
        return f"[Error processing MSG file: {e}. Is extract_msg installed?]"
    
def parse_txt(file_path: Union[str, Path]) -> str:
    """
    Parses a plain text file (UTF-8 encoding).

    Args:
        file_path: Path to the text file.

    Returns:
        The content of the file as a string.

    Raises:
        FileHandlingError: If the file is not found or cannot be read.
        ParsingError: If decoding fails or other unexpected errors occur.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse TXT file: {_file_path}")
    try:
        with open(_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        ASCIIColors.debug(f"Successfully parsed TXT file: {_file_path}")
        return content
    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except UnicodeDecodeError as e:
        msg = f"Encoding error parsing TXT file {_file_path} as UTF-8: {e}"
        ASCIIColors.error(msg)
        raise ParsingError(msg) from e
    except OSError as e:
        msg = f"OS error reading TXT file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e:
        msg = f"Unexpected error parsing TXT file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


def parse_pdf(file_path: Union[str, Path]) -> str:
    """
    Parses a PDF file to extract text content using pypdf.

    Args:
        file_path: Path to the PDF file.

    Returns:
        The extracted text content, concatenated from all pages.

    Raises:
        ConfigurationError: If 'pypdf' is not installed.
        FileHandlingError: If the file is not found or is empty.
        ParsingError: If the PDF structure is invalid or text extraction fails.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse PDF file: {_file_path}")
    try:
        # Import dynamically to check for dependency
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError, EmptyFileError
    except ImportError as e:
        msg = "Parsing PDF files requires 'pypdf'. Install with: pip install safe_store[parsing]"
        ASCIIColors.error(msg)
        raise ConfigurationError(msg) from e

    full_text = ""
    try:
        # Use strict=False to be more tolerant of minor PDF spec violations
        reader: PdfReaderProtocol = PdfReader(_file_path, strict=False)
        num_pages = len(reader.pages)
        ASCIIColors.debug(f"PDF '{_file_path.name}' loaded with {num_pages} pages.")

        if num_pages == 0:
             ASCIIColors.warning(f"PDF file '{_file_path.name}' contains zero pages.")
             return "" # Return empty string for zero-page PDFs

        for i, page in enumerate(reader.pages):
            page_num = i + 1
            try:
                page_text = page.extract_text()
                if page_text:
                    full_text += page_text + "\n" # Add newline between pages
                else:
                    ASCIIColors.debug(f"No text extracted from page {page_num} of '{_file_path.name}'.")
            except Exception as page_err: # Catch errors during individual page extraction
                 ASCIIColors.warning(f"Could not extract text from page {page_num} of '{_file_path.name}': {page_err}")
                 # Decide whether to continue or fail? Continue for now.
                 # Consider adding a flag to control strictness on page errors.

        ASCIIColors.debug(f"Successfully parsed PDF file: {_file_path}")
        return full_text.strip() # Remove leading/trailing whitespace

    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except EmptyFileError as e:
         msg = f"Cannot parse empty PDF file: {_file_path}"
         ASCIIColors.error(msg)
         # Consider empty file a FileHandlingError subtype? Or keep as ParsingError?
         # Let's use FileHandlingError as it relates to the file state.
         raise FileHandlingError(msg) from e
    except PdfReadError as e: # Catch specific pypdf structural errors
        msg = f"Error reading PDF structure in {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise ParsingError(msg) from e
    except OSError as e: # Catch potential OS errors during file reading by pypdf
        msg = f"OS error reading PDF file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e: # Catch other unexpected errors from pypdf
        msg = f"Unexpected error parsing PDF file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_pptx(file_path: Union[str, Path]) -> str:
    """
    Parses a PPTX file to extract text content using python-pptx.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        The extracted text content, concatenated from all slides and shapes.

    Raises:
        ConfigurationError: If 'python-pptx' is not installed.
        FileHandlingError: If the file is not found or OS error occurs during reading.
        ParsingError: If the file is not a valid PPTX format or text extraction fails.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse PPTX file: {_file_path}")

    try:
        # Import dynamically
        from pptx import Presentation
    except ImportError as e:
        pm.ensure_packages("python-pptx")
        try:
            # Import dynamically
            from pptx import Presentation
        except ImportError as e:
            msg = "Parsing PPTX files requires 'python-pptx'. Install with: pip install python-pptx"
            ASCIIColors.error(msg)
            raise ConfigurationError(msg) from e

    full_text = ""
    try:
        presentation = Presentation(_file_path)
        ASCIIColors.debug(f"PPTX '{_file_path.name}' loaded.")

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"

        ASCIIColors.debug(f"Successfully parsed PPTX file: {_file_path}")
        return full_text.strip()  # Remove leading/trailing whitespace

    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except OSError as e:  # Catch potential OS errors during file reading by python-pptx
        msg = f"OS error reading PPTX file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e:  # Catch other unexpected errors from python-pptx
        msg = f"Unexpected error parsing PPTX file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


def parse_docx(file_path: Union[str, Path]) -> str:
    """
    Parses a DOCX file to extract text content using python-docx.

    Args:
        file_path: Path to the DOCX file.

    Returns:
        The extracted text content, concatenated from all paragraphs.

    Raises:
        ConfigurationError: If 'python-docx' is not installed.
        FileHandlingError: If the file is not found.
        ParsingError: If the file is not a valid DOCX format or text extraction fails.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse DOCX file: {_file_path}")
    try:
        # Import dynamically
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as e:
        import pipmaster as pm
        pm.ensure_packages("python-docx")
        try:
            # Import dynamically
            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError
        except ImportError as e:
            msg = "Parsing DOCX files requires 'python-docx'. Install with: pip install safe_store[parsing]"
            ASCIIColors.error(msg)
            raise ConfigurationError(msg) from e

    full_text = ""
    try:
        document: DocumentProtocol = Document(_file_path)
        ASCIIColors.debug(f"DOCX '{_file_path.name}' loaded.")
        for para in document.paragraphs:
            full_text += para.text + "\n" # Add newline between paragraphs
        ASCIIColors.debug(f"Successfully parsed DOCX file: {_file_path}")
        return full_text.strip() # Remove leading/trailing whitespace

    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except PackageNotFoundError as e: # Specific error for invalid zip/docx format
         msg = f"File is not a valid DOCX (Zip) file: {_file_path}. Error: {e}"
         ASCIIColors.error(msg)
         raise ParsingError(msg) from e
    except OSError as e: # Catch potential OS errors during file reading by python-docx
        msg = f"OS error reading DOCX file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e: # Catch other unexpected errors from python-docx
        msg = f"Unexpected error parsing DOCX file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


def parse_html(file_path: Union[str, Path]) -> str:
    """
    Parses an HTML file to extract text content using BeautifulSoup.

    Uses 'lxml' if available for performance, otherwise falls back to
    Python's built-in 'html.parser'.

    Args:
        file_path: Path to the HTML file.

    Returns:
        The extracted text content, with tags stripped.

    Raises:
        ConfigurationError: If 'beautifulsoup4' (or 'lxml' if used) is not installed.
        FileHandlingError: If the file is not found or cannot be read.
        ParsingError: If the file cannot be decoded or parsed as HTML.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse HTML file: {_file_path}")
    try:
        # Import dynamically
        from bs4 import BeautifulSoup
        try:
            import lxml # noqa F401 ensure it's installed for best performance
            HTML_PARSER = 'lxml'
            ASCIIColors.debug("Using 'lxml' for HTML parsing.")
        except ImportError:
            ASCIIColors.debug("lxml not found, using Python's built-in 'html.parser' for HTML.")
            HTML_PARSER = 'html.parser'
    except ImportError as e:
        # This catches missing BeautifulSoup4
        msg = "Parsing HTML files requires 'BeautifulSoup4'. Install with: pip install safe_store[parsing]"
        ASCIIColors.error(msg)
        raise ConfigurationError(msg) from e

    try:
        with open(_file_path, 'r', encoding='utf-8') as f:
            # Read content first, then parse
             content = f.read()

        # Parse the HTML content
        soup: BeautifulSoupProtocol = BeautifulSoup(content, HTML_PARSER)
        # Get text, stripping extra whitespace and using newline as separator
        text = soup.get_text(separator='\n', strip=True)
        ASCIIColors.debug(f"Successfully parsed HTML file: {_file_path}")
        return text

    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except UnicodeDecodeError as e:
        msg = f"Encoding error parsing HTML file {_file_path} as UTF-8: {e}"
        ASCIIColors.error(msg)
        raise ParsingError(msg) from e
    except OSError as e:
        msg = f"OS error reading HTML file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e: # Catch potential BeautifulSoup errors or other issues
        msg = f"Error parsing HTML file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


# --- Dispatcher Function ---

# Define a type hint for the parser functions
ParserFunc = Callable[[Union[str, Path]], str]

# Map extensions to parser functions
parser_map: Dict[str, ParserFunc] = {
    '.txt': parse_txt,
    '.pdf': parse_pdf,
    '.docx': parse_docx,
    '.html': parse_html,
    '.htm': parse_html, # Treat .htm the same as .html
    '.pptx': parse_pptx,
    
    # Outlook messages
    '.msg': parse_msg,

    # --- General Text & Document Formats ---
    '.md': parse_txt,        # Markdown
    '.rst': parse_txt,       # reStructuredText
    '.tex': parse_txt,       # LaTeX source
    '.rtf': parse_txt,       # Rich Text Format (basic text extraction)
    '.log': parse_txt,       # Log files
    '.text': parse_txt,      # Generic text
    '.me': parse_txt,        # Often README files
    '.org': parse_txt,       # Emacs Org-mode

    # --- Data Serialization Formats (parsed as raw text) ---
    # If structured parsing is needed, dedicated parsers would be better.
    # For now, we treat them as text sources.
    '.json': parse_txt,      # JavaScript Object Notation
    '.xml': parse_txt,       # Extensible Markup Language
    '.csv': parse_txt,       # Comma-Separated Values
    '.tsv': parse_txt,       # Tab-Separated Values
    '.yaml': parse_txt,      # YAML Ain't Markup Language
    '.yml': parse_txt,       # YAML Ain't Markup Language (alternative)
    '.sql': parse_txt,       # SQL scripts
    '.graphql': parse_txt,   # GraphQL query language
    '.gql': parse_txt,       # GraphQL query language (alternative)

    # --- Configuration Files ---
    '.ini': parse_txt,       # INI configuration
    '.cfg': parse_txt,       # Configuration file
    '.conf': parse_txt,      # Configuration file
    '.toml': parse_txt,      # Tom's Obvious, Minimal Language
    '.env': parse_txt,       # Environment variables (e.g., .env files)
    '.properties': parse_txt,# Java properties files
    '. Htaccess': parse_txt, # Apache configuration (note: often no prefix dot in map)
    '.dockerfile': parse_txt,# Docker configuration
    'dockerfile': parse_txt, # Docker configuration (common naming)
    '.gitattributes': parse_txt, # Git attributes
    '.gitconfig': parse_txt,   # Git configuration
    '.gitignore': parse_txt,   # Git ignore patterns

    # --- Programming Language Source Code ---
    # Scripting Languages
    '.py': parse_txt,        # Python
    '.pyw': parse_txt,       # Python (GUI, no console)
    '.js': parse_txt,        # JavaScript
    '.mjs': parse_txt,       # ECMAScript modules
    '.cjs': parse_txt,       # CommonJS modules
    '.ts': parse_txt,        # TypeScript
    '.tsx': parse_txt,       # TypeScript with JSX (React)
    '.rb': parse_txt,        # Ruby
    '.pl': parse_txt,        # Perl
    '.pm': parse_txt,        # Perl Module
    '.php': parse_txt,       # PHP
    '.phtml': parse_txt,     # PHP templating
    '.sh': parse_txt,        # Shell script (Bourne, Bash, etc.)
    '.bash': parse_txt,      # Bash script
    '.zsh': parse_txt,       # Zsh script
    '.fish': parse_txt,      # Fish shell script
    '.ps1': parse_txt,       # PowerShell
    '.psm1': parse_txt,      # PowerShell Module
    '.psd1': parse_txt,      # PowerShell Data File
    '.lua': parse_txt,       # Lua
    '.r': parse_txt,         # R language
    '.tcl': parse_txt,       # TCL
    '.dart': parse_txt,      # Dart

    # Compiled Languages (Source)
    '.c': parse_txt,         # C
    '.h': parse_txt,         # C/C++ Header
    '.cpp': parse_txt,       # C++
    '.cxx': parse_txt,       # C++ (alternative)
    '.cc': parse_txt,        # C++ (alternative)
    '.hpp': parse_txt,       # C++ Header
    '.hxx': parse_txt,       # C++ Header (alternative)
    '.cs': parse_txt,        # C#
    '.java': parse_txt,      # Java
    '.go': parse_txt,        # Go
    '.rs': parse_txt,        # Rust
    '.swift': parse_txt,     # Swift
    '.kt': parse_txt,        # Kotlin
    '.kts': parse_txt,       # Kotlin Script
    '.scala': parse_txt,     # Scala
    '.m': parse_txt,         # Objective-C or MATLAB (source is text)
    '.mm': parse_txt,        # Objective-C++
    '.f': parse_txt,         # Fortran
    '.for': parse_txt,       # Fortran
    '.f90': parse_txt,       # Fortran 90
    '.f95': parse_txt,       # Fortran 95
    '.pas': parse_txt,       # Pascal
    '.d': parse_txt,         # D language
    '.vb': parse_txt,        # Visual Basic .NET
    '.vbs': parse_txt,       # VBScript

    # Web Development & Templating (beyond just HTML)
    '.css': parse_txt,       # Cascading Style Sheets
    '.scss': parse_txt,      # Sass CSS preprocessor
    '.sass': parse_txt,      # Sass CSS preprocessor (indented syntax)
    '.less': parse_txt,      # Less CSS preprocessor
    '.styl': parse_txt,      # Stylus CSS preprocessor
    '.jsx': parse_txt,       # JavaScript XML (React)
    '.vue': parse_txt,       # Vue.js single-file components
    '.svelte': parse_txt,    # Svelte components
    '.ejs': parse_txt,       # Embedded JavaScript templates
    '.hbs': parse_txt,       # Handlebars templates
    '.mustache': parse_txt,  # Mustache templates
    '.jinja': parse_txt,     # Jinja templates (Python)
    '.jinja2': parse_txt,    # Jinja2 templates
    '.twig': parse_txt,      # Twig templates (PHP)
    '.erb': parse_txt,       # Embedded Ruby (Rails templates)
    '.jsp': parse_txt,       # JavaServer Pages
    '.asp': parse_txt,       # Active Server Pages (classic)
    '.aspx': parse_txt,      # ASP.NET Web Forms

    # Other/Assembly/Specialized
    '.asm': parse_txt,       # Assembly language
    '.s': parse_txt,         # Assembly language (common on Unix)
    '.bat': parse_txt,       # Batch file (Windows)
    '.cmd': parse_txt,       # Command script (Windows NT)
    '.vhd': parse_txt,       # VHDL (Hardware description language)
    '.vhdl': parse_txt,      # VHDL
    '.sv': parse_txt,        # SystemVerilog
    '.bib': parse_txt,       # BibTeX bibliography file
    '.srt': parse_txt,       # SubRip Subtitle file
    '.sub': parse_txt,       # Subtitle file
    '.vtt': parse_txt,       # WebVTT Subtitle file
    '.po': parse_txt,        # Portable Object (localization)
    '.pot': parse_txt,       # Portable Object Template (localization)
    '.strings': parse_txt,   # iOS/macOS localization strings        
}

SAFE_STORE_SUPPORTED_FILE_EXTENSIONS = list(parser_map.keys())

def parse_document(file_path: Union[str, Path]) -> str:
    """
    Parses a document based on its file extension.

    Dispatches to the appropriate parser (.txt, .pdf, .docx, .html/.htm).

    Args:
        file_path: Path to the document file.

    Returns:
        The extracted text content of the document.

    Raises:
        FileHandlingError: If the input path is not a file.
        ConfigurationError: If the file extension is unsupported.
        ParsingError: If the dispatched parser encounters an error.
    """
    _file_path = Path(file_path)
    if not _file_path.is_file():
        msg = f"Input path is not a file: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) # Use specific error

    extension = _file_path.suffix.lower()
    ASCIIColors.debug(f"Dispatching parser for extension '{extension}' on file: {_file_path.name}")



    parser_func = parser_map.get(extension)

    if parser_func:
        try:
            # Call the appropriate parser function
            return parser_func(_file_path)
        except (ConfigurationError, FileHandlingError, ParsingError) as e:
             # Re-raise specific known errors directly
             raise e
        except Exception as e:
             # Wrap unexpected errors from the specific parser
             msg = f"Unexpected error during parsing dispatch for {_file_path.name} (extension '{extension}'): {e}"
             ASCIIColors.error(msg, exc_info=True)
             # Wrap in ParsingError as it occurred during the parsing stage
             raise ParsingError(msg) from e
    else:
        # Handle unsupported file type
        msg = f"Unsupported file type extension: '{extension}' for file: {_file_path}. No parser available."
        ASCIIColors.warning(msg)
        # Use ConfigurationError for unsupported type, as it's a setup/config issue
        raise ConfigurationError(msg)
