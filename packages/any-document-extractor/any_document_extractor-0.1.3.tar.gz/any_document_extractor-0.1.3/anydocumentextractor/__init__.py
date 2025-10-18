import os

from docx import Document
from openpyxl import load_workbook
from pdfminer.high_level import extract_text
from pptx import Presentation

class DocumentExtractor:
    def __init__(self, fp: str):
        self.fp = fp
        self.extractors = {
            '.docx': self.extract_word,
            '.xlsx': self.extract_excel,
            '.pdf': self.extract_pdf,
            '.pptx': self.extract_ppt
        }

    def _remove_all_whitespace_translate(self, text):
        whitespace_chars = ' \t\n\r\x0b\x0c'
        trans_table = str.maketrans('', '', whitespace_chars)
        return text.translate(trans_table)

    def extract_word(self):
        doc = Document(self.fp)
        data =  ''.join([para.text for para in doc.paragraphs])
        return self._remove_all_whitespace_translate(data)

    def extract_excel(self):
        wb = load_workbook(self.fp, data_only=True)
        ws = wb.active
        data = [[cell.value for cell in row] for row in ws.iter_rows()]
        table_str = ""
        for row in data:
            row_str = "".join(str(cell) if cell is not None else "" for cell in row)
            table_str += row_str

        return self._remove_all_whitespace_translate(table_str.strip())

    def extract_pdf(self):
        text = extract_text(self.fp)
        return self._remove_all_whitespace_translate(text)

    def extract_ppt(self):
        prs = Presentation(self.fp)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    text_content.append(shape.text_frame.text)
        data = ''.join(text_content)
        return self._remove_all_whitespace_translate(data)

    def extract(self):
        ext = os.path.splitext(self.fp)[1].lower()
        if ext not in self.extractors:
            raise ValueError(f"不支持: {ext}")
        return self.extractors[ext]()